"""Drive Files REST API (PR2.5)

端点:
  POST   /api/v1/drive/files          → 创建文件元数据 (multipart 完成后调)
  GET    /api/v1/drive/files          → 列 drive 文件 (含分页/filter)
  GET    /api/v1/drive/files/{id}     → 详情
  PUT    /api/v1/drive/files/{id}     → 改名/移动/改 visibility
  DELETE /api/v1/drive/files/{id}     → 软删
  POST   /api/v1/drive/files/{id}/restore → 恢复 (3 天保留期内)
  POST   /api/v1/drive/files/{id}/extract-to-kb → 升级到公共知识库
  GET    /api/v1/drive/storage-stats  → 容量统计 (per member)
  GET    /api/v1/drive/files/{id}/preview → 预览元信息 (image/pdf/text, 60s Redis 缓存) [v2 PR8e]
  GET    /api/v1/drive/files/{id}/thumbnail → 缩略图 URL

Multipart 简化 (PR2.3): 单端点流式接收 + minio 自管分片
  POST   /api/v1/drive/files/upload   → multipart 接收 (单端点, content-type: multipart/form-data)
                                          form: {file, filename, content_type, total_size,
                                                 folder_id?, visibility?, storage_mode?}
"""
import io
import logging
import math
from datetime import datetime, timedelta
from typing import List, Optional

from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, Request, UploadFile
from pydantic import BaseModel, Field
from sqlalchemy import select, text
from sqlalchemy.ext.asyncio import AsyncSession

from app.api.v1._drive_error_helper import (
    ERR_BATCH_PARAM_MISSING,
    ERR_FILE_FORBIDDEN,
    ERR_FILE_NOT_FOUND,
    ERR_SHARE_LINK_NOT_FOUND,
    raise_app_error,
)
from app.core.database import get_db
from app.core.security import get_current_user
from app.models.knowledge import Knowledge
from app.models.folder import Folder
from app.models.member import Member
from app.services.drive_service import (
    DriveService,
    DriveServiceError,
    MAX_DRIVE_FILE_SIZE_BYTES,
    MAX_DRIVE_FILE_SIZE_MB,
)
from app.services.file_service import file_service
from app.services.generic_chunked_upload_service import (
    ChunkedUploadError,
    generic_chunked_upload_service,
)
from app.services.folder_service import FolderService

logger = logging.getLogger("microbubble.drive_api")
router = APIRouter(prefix="/drive", tags=["网盘文件"])


# === Schemas ===

class DriveFileItem(BaseModel):
    id: int
    title: str
    file_path: str
    file_name: str
    file_type: str
    file_size: int
    storage_mode: str
    visibility: str
    folder_id: Optional[int] = None
    created_by: Optional[int] = None
    # 批次⑩: 上传者三件套 — 之前 schema 漏声明, pydantic 静默丢弃导致上传者列恒 '—'
    owner_name: Optional[str] = None
    owner_username: Optional[str] = None
    owner_avatar: Optional[str] = None
    source_type: Optional[str] = None
    created_at: Optional[str] = None
    updated_at: Optional[str] = None
    deleted_at: Optional[str] = None
    original_parent_id: Optional[int] = None
    original_path: Optional[str] = None
    remaining_days: Optional[int] = None
    auto_delete_at: Optional[str] = None
    download_count: int = 0
    share_token: Optional[str] = None
    share_expires_at: Optional[str] = None
    # v2 PR2: 收藏字段
    is_starred: bool = False
    starred_at: Optional[str] = None
    # 批次① B6 搜索结果配套: 所属文件夹名 (列表端点批量 SELECT, 搜索态必带; 前端
    # FileCard/三栏工作台"位置"列消费)。folder_id 为 None (顶级) 时也返回 None。
    folder_name: Optional[str] = None
    # v2 PR4: 秒传 + 版本历史
    file_hash: Optional[str] = None
    is_latest: bool = True
    version_number: int = 1
    # v2 PR5: 缩略图字段
    thumbnail_path: Optional[str] = None
    thumbnail_status: str = "pending"  # pending | ready | failed
    # v2 PR6-P19: 团队共享盘隔离 (uploaded from specialView='team' = True)
    # 决定 list_drive_files view=personal|team 过滤
    is_team_shared: bool = False

    class Config:
        from_attributes = True


class DriveFileListResponse(BaseModel):
    items: List[DriveFileItem]
    total: int
    page: int
    page_size: int


class DriveFileUpdate(BaseModel):
    title: Optional[str] = Field(None, max_length=500)
    file_name: Optional[str] = Field(None, max_length=500)  # PR4.4: 修复 PR2.5 漏的 file_name 字段
    visibility: Optional[str] = None  # private | team | public
    folder_id: Optional[int] = None  # 0 = move to root


class ExtractToKBRequest(BaseModel):
    target_visibility: str = "team"  # team | public (不能 private)


class StorageStatsResponse(BaseModel):
    file_count: int
    by_visibility: dict
    active: bool = True


# === v2 PR5 Schemas: 配额 + 分片 + 缩略图 ===

class StorageQuotaResponse(BaseModel):
    """GET /api/v1/drive/storage-quota 响应"""
    user_id: int
    used_bytes: int
    quota_bytes: int
    percent: float
    file_count: int
    is_over_quota: bool
    updated_at: Optional[str] = None


class ChunkedUploadInitRequest(BaseModel):
    """POST /api/v1/drive/files/upload/init 请求"""
    file_name: str = Field(..., max_length=500)
    file_size: int = Field(..., gt=0, le=2 * 1024 * 1024 * 1024)  # 上限 2GB
    total_chunks: int = Field(..., ge=1, le=2000)
    file_hash: Optional[str] = Field(None, max_length=64)
    folder_id: Optional[int] = None
    visibility: str = "team"
    # v2 PR6-P19: is_team_shared 在 complete 阶段才传, 这里不接 (前端可在最后决定)


class ChunkedUploadInitResponse(BaseModel):
    """POST /api/v1/drive/files/upload/init 响应"""
    upload_id: str  # session_id
    object_name: str  # 临时占位
    total_chunks: int
    chunk_size_hint: int = 5 * 1024 * 1024  # 5MB 提示
    uploaded_chunks: List[int] = []
    expires_at: str


class ChunkedUploadStatusResponse(BaseModel):
    """GET /api/v1/drive/files/upload/{id} 响应 (断点续传)"""
    upload_id: str
    file_name: str
    file_size: int
    total_chunks: int
    uploaded_chunks: List[int]
    status: str  # active | completed | aborted
    expires_at: str


class ChunkedUploadCompleteRequest(BaseModel):
    """POST /api/v1/drive/files/upload/{id}/complete 请求"""
    change_note: Optional[str] = Field(None, max_length=500)
    # v2 PR6-P19: 团队共享盘标识 (init 时已传过, complete 时可再传覆盖)
    is_team_shared: Optional[bool] = Field(
        None, description="v2 PR6-P19: 覆盖 init 时的设置 (None=沿用)"
    )


class ThumbnailResponse(BaseModel):
    """GET /api/v1/drive/files/{id}/thumbnail 响应 (返 URL)"""
    file_id: int
    thumbnail_path: Optional[str] = None
    thumbnail_status: str  # pending | ready | failed
    thumbnail_url: Optional[str] = None  # MinIO 公开读 URL 或 None


class PreviewResponse(BaseModel):
    """v2 PR8e: GET /api/v1/drive/files/{id}/preview 响应 (轻量元信息, 60s Redis 缓存)

    用途: 前端 grid/list 卡片快速拿到预览类型 + 缩略图 URL + 文字预览头部 + PDF 页数,
    避免每次点开都拉整文件 + 让 file-detail 页能秒开"封面信息"块.

    字段:
    - file_id / file_name / file_type / file_size: 复用 DriveFileItem 核心字段
    - preview_type: image | pdf | text | other (前端按此选择渲染策略)
    - cached: True=从 Redis 缓存返, False=实时解析
    - width / height: 仅 image 类型, 用 Pillow 取真实像素
    - page_count: 仅 pdf 类型, 用 PyMuPDF 取
    - thumbnail_url: 复用 thumbnail_path 的 MinIO 公开 URL (无 JWT, 浏览器直拉)
    - text_preview: 仅 text 类型, 前 1KB UTF-8 字符串 (空字符串表示 binary)
    - first_page_url: 仅 pdf 类型, 公开读 URL (无 ?disposition=inline, 让浏览器用 PDF viewer 渲染)
    """
    file_id: int
    file_name: str
    file_type: str
    file_size: int
    preview_type: str  # image | pdf | text | other
    cached: bool = False
    width: Optional[int] = None
    height: Optional[int] = None
    page_count: Optional[int] = None
    thumbnail_url: Optional[str] = None
    text_preview: Optional[str] = None  # 最多 1KB, text-only
    first_page_url: Optional[str] = None  # pdf 公开读 (CDN/预签名)


async def _owner_lookup(db: AsyncSession, items) -> dict:
    """批次⑦: 批量查本页文件 created_by 的 Member — 之前只 list 端点没传 owner_lookup,
    上传者列恒 '—'。返回 {member_id: Member}, 空页短路。"""
    uids = {x.created_by for x in items if x.created_by is not None}
    if not uids:
        return {}
    rows = (await db.execute(select(Member).where(Member.id.in_(uids)))).scalars().all()
    return {m.id: m for m in rows}


async def _folder_name_map(db: AsyncSession, items) -> dict:
    """批量查本页文件所属 folder {id: name} — B6 搜索结果 / 三栏工作台"位置"列, 防 N+1。"""
    fids = {x.folder_id for x in items if x.folder_id is not None}
    if not fids:
        return {}
    rows = (await db.execute(
        select(Folder.id, Folder.name).where(Folder.id.in_(fids))
    )).all()
    return {r[0]: r[1] for r in rows}


def _to_item(
    k: Knowledge,
    owner_lookup: Optional[dict] = None,
    starred_ids: Optional[set] = None,
    folder_map: Optional[dict] = None,
) -> DriveFileItem:
    """Build DriveFileItem.

    v2.16 (2026-07-11): 可选 owner_lookup 传 {member_id: Member} 字典,
    一次 JOIN 查所有 owner 后批量 attach (避免 N+1 query).
    老 callsite 不传 owner_lookup 时, owner_name/owner_username 字段为 None,
    前端 fallback 到 "#<created_by>" 显示 (CLAUDE.md 兼容原则).

    批次① 收藏个人化 (alembic 134): 可选 starred_ids 传"当前成员收藏过"的
    {file_id} 集合 (列表端点先批量 SELECT), **响应字段名 is_starred 不变**但语义
    升级为 per-user 视图; 不传时 fallback legacy 列 (详情/老 callsite 兼容)。
    """
    owner_name = None
    owner_username = None
    owner_avatar = None
    if owner_lookup and k.created_by in owner_lookup:
        m = owner_lookup[k.created_by]
        owner_name = m.name if m.name else None
        owner_username = m.username if m.username else None
        owner_avatar = m.avatar if m.avatar else None

    folder_name = None
    if folder_map and k.folder_id is not None:
        folder_name = folder_map.get(k.folder_id)

    remaining_days = None
    auto_delete_at = None
    if k.deleted_at:
        from app.config import settings

        retention_days = getattr(settings, "DRIVE_RETENTION_DAYS", 30)
        expires_at = k.deleted_at + timedelta(days=retention_days)
        seconds_left = (expires_at - datetime.utcnow()).total_seconds()
        remaining_days = max(0, math.ceil(seconds_left / 86400))
        auto_delete_at = expires_at.isoformat()

    return DriveFileItem(
        id=k.id,
        title=k.title,
        file_path=k.file_path or "",
        file_name=k.file_name or "",
        file_type=k.file_type or "",
        file_size=k.file_size or 0,  # PR4: 真值 (PR2 之前 0)
        storage_mode=k.storage_mode,
        visibility=k.visibility,
        folder_id=k.folder_id,
        created_by=k.created_by,
        owner_name=owner_name,
        owner_username=owner_username,
        owner_avatar=owner_avatar,
        source_type=k.source_type,
        created_at=str(k.created_at) if k.created_at else None,
        updated_at=str(k.updated_at) if k.updated_at else None,
        deleted_at=str(k.deleted_at) if k.deleted_at else None,
        original_parent_id=k.original_parent_id,
        original_path=k.original_path,
        remaining_days=remaining_days,
        auto_delete_at=auto_delete_at,
        download_count=k.download_count or 0,
        share_token=k.share_token,
        share_expires_at=str(k.share_expires_at) if k.share_expires_at else None,
        # 收藏个人化: 传了 starred_ids 即 per-user 视图; legacy starred_at 仅在
        # 无 per-user 覆盖时透出 (前端主要消费 is_starred)
        is_starred=(k.id in starred_ids) if starred_ids is not None else bool(k.is_starred),
        starred_at=str(k.starred_at) if k.starred_at else None,
        folder_name=folder_name,
        file_hash=k.file_hash,        # PR4
        is_latest=bool(k.is_latest),  # PR4
        version_number=k.version_number or 1,  # PR4
        thumbnail_path=k.thumbnail_path,  # PR5
        thumbnail_status=k.thumbnail_status or "pending",  # PR5
        is_team_shared=bool(k.is_team_shared),  # v2 PR6-P19
    )


# === 单端点 multipart 上传 (PR2.3 简化版) ===

@router.post("/files/upload", response_model=DriveFileItem, status_code=201)
async def upload_drive_file(
    file: UploadFile = File(..., description="文件二进制"),
    filename: Optional[str] = Form(None, description="原始文件名 (默认 file.filename)"),
    content_type: Optional[str] = Form(None, description="MIME (默认 file.content_type)"),
    folder_id: Optional[int] = Form(None, description="目标 folder id (None=顶级)"),
    visibility: str = Form("team", description="private|team|public"),
    # v2 PR6-P19: 用户上传时所在的视图 (specialView='team' = True)
    is_team_shared: bool = Form(False, description="v2 PR6-P19: True=团队共享盘上传, 不在个人网盘显示"),
    title: Optional[str] = Form(None, description="文件标题 (默认 = filename)"),
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """单端点 drive 文件上传 (FastAPI 接收 multipart, minio 自管分片)

    不超过 2GB (MAX_DRIVE_FILE_SIZE_BYTES); 走 init/complete pattern 简化为单端点.

    v2 PR6-P19: is_team_shared 标记上传来源, 前端在团队共享盘视图上传时 = true,
    个人视图上传 = false. 后端 list_drive_files 走 view=personal|team 隔离过滤.
    """
    real_filename = filename or file.filename or "unnamed"
    real_ct = content_type or file.content_type or "application/octet-stream"

    # 校验 visibility
    if visibility not in ("private", "team", "public"):
        raise HTTPException(status_code=400, detail=f"非法 visibility: {visibility}")

    # 校验 folder_id 存在 (2026-09 单一团队空间: 不再校验 folder owner, 任何成员可上传到任意 folder)
    if folder_id is not None:
        folder_svc = FolderService(db)
        folder = await folder_svc.get_folder(folder_id)
        if folder is None:
            raise HTTPException(status_code=404, detail=f"folder {folder_id} 不存在")

    # 读取 body
    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="上传文件为空")
    if len(data) > MAX_DRIVE_FILE_SIZE_BYTES:
        raise HTTPException(status_code=413, detail=f"文件超过 {MAX_DRIVE_FILE_SIZE_MB}MB")

    # 算 folder_path (e.g. "1/4/") from folder_id
    folder_path = None
    if folder_id is not None:
        folder_svc = FolderService(db)
        folder = await folder_svc.get_folder(folder_id)
        if folder is not None:
            folder_path = folder.path

    # 1) init
    try:
        init_resp = await generic_chunked_upload_service.init_upload(
            filename=real_filename,
            content_type=real_ct,
            total_size=len(data),
            folder_path=folder_path,
            storage_mode="drive",
            user_id=current_user.id,
        )
    except ChunkedUploadError as e:
        raise HTTPException(status_code=e.status_code, detail=str(e))

    # 2) complete (写 MinIO)
    try:
        complete_resp = await generic_chunked_upload_service.complete_upload(
            upload_id=init_resp["upload_id"],
            data=data,
            content_type=real_ct,
        )
    except ChunkedUploadError as e:
        # 失败回滚 init (清空)
        await generic_chunked_upload_service.abort_upload(
            upload_id=init_resp["upload_id"],
        )
        raise HTTPException(status_code=e.status_code, detail=str(e))

    # 3) 落 Knowledge 元数据
    drive_svc = DriveService(db)
    try:
        knowledge = await drive_svc.create_file(
            title=title or real_filename,
            file_path=complete_resp["object_name"],
            file_name=real_filename,
            file_type="." + real_filename.rsplit(".", 1)[-1] if "." in real_filename else "",
            file_size=complete_resp["size"] or len(data),
            owner_id=current_user.id,
            storage_mode="drive",
            visibility=visibility,
            folder_id=folder_id,
            created_by=current_user.id,
            source_type="drive",
            is_team_shared=is_team_shared,  # v2 PR6-P19 (2026-09 单一团队空间: 入参已忽略, create_file 服务端恒置 True)
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=str(e))

    # PR5: Fire-and-forget 缩略图生成 + 配额重算
    # 2026-09-05: + 网盘文件默认入库 RAG (上传即自动 drive → kb)
    try:
        from app.services.thumbnail_tasks import generate_thumbnail_task
        from app.services.storage_tasks import recalc_user_storage_task
        from app.config import settings as _settings

        generate_thumbnail_task.delay(knowledge.id)
        recalc_user_storage_task.delay(current_user.id)
        if _settings.DRIVE_AUTO_INGEST_KB:
            from app.services.drive_ingest_tasks import auto_ingest_drive_file_task
            auto_ingest_drive_file_task.delay(knowledge.id)
    except Exception as e:
        logger.warning(f"[drive.upload] fire Celery 失败 (非阻塞): {e}")

    logger.info(
        f"[drive.upload] id={knowledge.id} file_name={real_filename} "
        f"size={complete_resp['size']} visibility={visibility} folder_id={folder_id}"
    )
    return _to_item(knowledge)


# === CRUD ===

@router.get("/files", response_model=DriveFileListResponse)
async def list_drive_files(
    folder_id: Optional[int] = Query(None, description="父 folder id"),
    visibility: Optional[str] = Query(None, description="private|team|public"),
    page: int = Query(1, ge=1),
    page_size: int = Query(50, ge=1, le=100),
    include_deleted: bool = Query(False),
    # v2 PR2: sort + filter 新参数
    sort_by: Optional[str] = Query(
        None,
        description="排序字段: created_at | updated_at | file_name | starred_at",
    ),
    sort_order: Optional[str] = Query("desc", description="asc | desc"),
    starred_only: bool = Query(False, description="仅显示收藏"),
    file_type: Optional[str] = Query(
        None,
        description="类型过滤: pdf | image | video | audio | word | ppt | excel | text (无值=全部, office 仍兼容)",
    ),
    # v2 PR6-P19: 视图隔离 (personal=个人网盘默认, team=团队共享盘, all=全显)
    view: Optional[str] = Query(
        "personal",
        description="视图隔离: personal (默认, 不含 is_team_shared=true) | team (仅 is_team_shared=true) | all (不过滤)",
    ),
    # v2.21 (2026-07-11): 团队共享盘顶级列出整个团队空间 (root + sub folder)
    # 默认 False 保持 v2 PR3 行为 (folder_id=None → folder_id IS NULL)
    # 个人顶级 view=personal 不该传 True (会 leak team view 行为)
    include_subfolders: bool = Query(
        False,
        description="v2.21: folder_id=None 时是否包含 sub folder PPT (🌐 团队共享盘顶级 view 用)",
    ),
    # 批次① B6 (2026-09-05): 文件名/标题中缀搜索 (转义 % _ 字面量; 空/不传 = 行为与老版完全一致)
    search: Optional[str] = Query(
        None,
        max_length=100,
        description="按文件名/标题模糊搜索 (全盘跨文件夹, 搜索态忽略 folder 约束)",
    ),
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """列 drive 文件 (2026-09 单一团队空间: view 参数继续接受但语义统一)

    v2 PR2: 支持 sort_by/sort_order/starred_only/file_type.
    v2 PR6-P19: view 参数隔离个人/团队共享盘. **已退役**: personal/team/all 均返回
      全量团队文件 (is_team_shared_filter 恒 None), 老客户端不炸。
    v2.21: include_subfolders 继续接受; 单一空间下不再按 view 屏蔽。
    批次① B6: search 非空时 folder 约束失效 (folder_id 置 None + include_subfolders
      置 True → 不加 folder filter, 全盘跨夹命中; 文件名中缀走 134 的 trgm GIN 索引)。
    批次① 收藏个人化: is_starred 响应字段升级为**当前成员**视角 (批量 SELECT 本页
      drive_file_stars, 无 N+1)。
    """
    # 参数合法性照旧校验 (老客户端可能传非法值)
    if view not in ("personal", "team", "all"):
        raise HTTPException(
            status_code=400,
            detail=f"无效 view 参数: '{view}', 必须是 personal|team|all",
        )
    # v2 PR6-P19 收口: 恒 None = 不按 is_team_shared 过滤 (数据迁移 133 全部回填 true)
    is_team_shared_filter = None

    # B6 搜索态: 忽略 folder 约束 (全盘搜), 空 search 走原路径 0 行为变化
    searching = bool(search and search.strip())
    if searching:
        folder_id = None
        include_subfolders = True

    svc = DriveService(db)
    try:
        items, total = await svc.list_files(
            current_user_id=current_user.id,
            folder_id=folder_id,
            visibility_filter=visibility,
            storage_mode="drive",
            include_deleted=include_deleted,
            page=page,
            page_size=page_size,
            sort_by=sort_by or "created_at",
            sort_order=sort_order or "desc",
            starred_only=starred_only,
            file_type=file_type,
            is_team_shared=is_team_shared_filter,
            include_subfolders=include_subfolders,
            search=search,
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=str(e))
    # 收藏个人化: 一次批量 SELECT 本页 ids 的本人 star 集合, attach 到 is_starred
    starred_ids = await svc.get_starred_ids([x.id for x in items], current_user.id)
    # B6 搜索态: 结果跨夹, 批量带出所属文件夹名 (非搜索态同样填充, 前端选择消费)
    folder_map = await _folder_name_map(db, items)
    owner_map = await _owner_lookup(db, items)
    return DriveFileListResponse(
        items=[_to_item(x, starred_ids=starred_ids, folder_map=folder_map, owner_lookup=owner_map) for x in items],
        total=total,
        page=page,
        page_size=page_size,
    )


@router.get("/files/{file_id}", response_model=DriveFileItem)
async def get_drive_file(
    file_id: int,
    include_deleted: bool = Query(False),
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """drive 文件详情 + 越权检查"""
    svc = DriveService(db)
    f = await svc.get_file(file_id, current_user_id=current_user.id)
    if f is None:
        raise HTTPException(
            status_code=404,
            detail="file 不存在或无权访问",
        )
    return _to_item(f, owner_lookup=await _owner_lookup(db, [f]))


@router.put("/files/{file_id}", response_model=DriveFileItem)
async def update_drive_file(
    file_id: int,
    payload: DriveFileUpdate,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """更新 drive 文件 (rename / move / change visibility)"""
    svc = DriveService(db)
    try:
        f = await svc.update_file(
            file_id,
            current_user_id=current_user.id,
            title=payload.title,
            file_name=payload.file_name,  # PR4.4: 透传 file_name
            visibility=payload.visibility,
            folder_id=payload.folder_id,
        )
    except DriveServiceError as e:
        # W1 T1 migration: DriveServiceError → AppException envelope (helper)
        raise_app_error(e.status_code, ERR_FILE_FORBIDDEN, str(e), file_id=file_id)
    if f is None:
        raise_app_error(404, ERR_FILE_NOT_FOUND, "file 不存在", file_id=file_id)
    return _to_item(f)


@router.delete("/files/{file_id}", status_code=204)
async def delete_drive_file(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """软删 drive 文件 (2026-09 单一团队空间: 任何成员可删)"""
    svc = DriveService(db)
    ok = await svc.soft_delete_file(file_id, current_user_id=current_user.id)
    if not ok:
        raise HTTPException(status_code=404, detail="file 不存在")
    return


@router.post("/files/{file_id}/restore", response_model=DriveFileItem)
async def restore_drive_file(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """恢复软删 drive 文件 (30 天保留期内; 2026-09-05 角色扁平化: 任何成员可恢复)."""
    svc = DriveService(db)
    f = await svc.restore_file(
        file_id,
        current_user_id=current_user.id,
        is_admin=True,
    )
    if f is None:
        raise HTTPException(status_code=404, detail="file 不存在或无恢复权限")
    return _to_item(f)


@router.post("/files/{file_id}/extract-to-kb", response_model=DriveFileItem)
async def extract_to_kb(
    file_id: int,
    payload: ExtractToKBRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """drive 文件升级到公共知识库 (storage_mode: drive → kb, source_type: drive → drive_extracted)

    后续 PR3 会触发 LLM 提取 + embedding 异步任务
    """
    svc = DriveService(db)
    try:
        f = await svc.extract_to_kb(
            file_id,
            current_user_id=current_user.id,
            target_visibility=payload.target_visibility,
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=str(e))
    if f is None:
        raise HTTPException(status_code=404, detail="file 不存在")
    return _to_item(f)


# ==========================================================
# v2 PR8.6: 文件级软锁 (防冲突协作)
# Redis TTL 5 分钟自动过期, 释放/重新获取强制覆盖
# ==========================================================

import time as _time

DRIVE_FILE_LOCK_TTL_SECONDS = 300  # 5 分钟
_DRIVE_LOCK_KEY_PREFIX = "drive:file_lock:"


async def _get_redis_client():
    """延迟创建 Redis 客户端 (CLAUDE.md 752 行铁律: 不在模块顶部创建)"""
    import redis.asyncio as aioredis
    from app.config import settings
    return aioredis.from_url(settings.REDIS_URL, decode_responses=True)


async def _drive_lock_get_holder(file_id: int) -> Optional[dict]:
    """读 Redis 锁 holder, 无锁返 None"""
    client = await _get_redis_client()
    try:
        raw = await client.get(f"{_DRIVE_LOCK_KEY_PREFIX}{file_id}")
        if not raw:
            return None
        import json
        try:
            data = json.loads(raw)
            # 校验 TTL 实际剩余 (Redis 5 分钟已重新校准)
            ttl = await client.ttl(f"{_DRIVE_LOCK_KEY_PREFIX}{file_id}")
            data["ttl_remaining"] = ttl if ttl > 0 else 0
            return data
        except Exception:
            return None
    finally:
        await client.aclose()


class FileLockResponse(BaseModel):
    """v2 PR8.6: 锁状态响应"""
    file_id: int
    locked: bool
    holder_user_id: Optional[int] = None
    holder_username: Optional[str] = None
    holder_name: Optional[str] = None
    acquired_at: Optional[str] = None
    ttl_remaining: int = 0  # 秒


@router.post("/files/{file_id}/lock", response_model=FileLockResponse)
async def acquire_file_lock(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """v2 PR8.6: 软锁 drive 文件 (5 分钟 Redis TTL)

    - 同用户重复 acquire: 续期 TTL (返回新 acquired_at)
    - 异用户 acquire: 返 409 + 当前 holder 信息
    """
    svc = DriveService(db)
    f = await svc.get_file(file_id, current_user_id=current_user.id)
    if f is None:
        raise HTTPException(status_code=404, detail="file 不存在或无权访问")

    client = await _get_redis_client()
    try:
        key = f"{_DRIVE_LOCK_KEY_PREFIX}{file_id}"
        import json
        now_iso = _time.strftime("%Y-%m-%dT%H:%M:%SZ", _time.gmtime())
        payload = json.dumps({
            "user_id": current_user.id,
            "username": getattr(current_user, "username", None) or "",
            "name": getattr(current_user, "name", "") or "",
            "acquired_at": now_iso,
        })
        existing_raw = await client.get(key)
        if existing_raw:
            try:
                existing = json.loads(existing_raw)
                if int(existing.get("user_id", 0)) != current_user.id:
                    return FileLockResponse(
                        file_id=file_id,
                        locked=True,
                        holder_user_id=int(existing.get("user_id", 0)) or None,
                        holder_username=existing.get("username"),
                        holder_name=existing.get("name"),
                        acquired_at=existing.get("acquired_at"),
                        ttl_remaining=max(await client.ttl(key), 0),
                    )
            except Exception:
                # 损坏的 JSON 直接覆盖
                pass
        # 同用户续期 / 异用户第一次拿到 → 写入
        await client.set(key, payload, ex=DRIVE_FILE_LOCK_TTL_SECONDS)
        return FileLockResponse(
            file_id=file_id,
            locked=True,
            holder_user_id=current_user.id,
            holder_username=getattr(current_user, "username", None) or "",
            holder_name=getattr(current_user, "name", "") or "",
            acquired_at=now_iso,
            ttl_remaining=DRIVE_FILE_LOCK_TTL_SECONDS,
        )
    finally:
        await client.aclose()


@router.delete("/files/{file_id}/lock", response_model=FileLockResponse)
async def release_file_lock(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """v2 PR8.6: 释放 drive 文件锁

    - 仅 holder 可释放, 异用户返 403
    - 锁不存在返 200 + locked=False (幂等)
    """
    svc = DriveService(db)
    f = await svc.get_file(file_id, current_user_id=current_user.id)
    if f is None:
        raise HTTPException(status_code=404, detail="file 不存在或无权访问")

    client = await _get_redis_client()
    try:
        key = f"{_DRIVE_LOCK_KEY_PREFIX}{file_id}"
        raw = await client.get(key)
        if not raw:
            return FileLockResponse(file_id=file_id, locked=False)
        import json
        try:
            existing = json.loads(raw)
            if int(existing.get("user_id", 0)) != current_user.id:
                raise HTTPException(status_code=403, detail="非锁持有者, 无法释放")
        except HTTPException:
            raise
        except Exception:
            pass
        await client.delete(key)
        return FileLockResponse(file_id=file_id, locked=False)
    finally:
        await client.aclose()


@router.get("/files/{file_id}/lock", response_model=FileLockResponse)
async def get_file_lock(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """v2 PR8.6: 查 drive 文件当前锁状态

    任何可见该文件的人都能查 (含 public/team), 无锁时 locked=False
    """
    svc = DriveService(db)
    f = await svc.get_file(file_id, current_user_id=current_user.id)
    if f is None:
        raise HTTPException(status_code=404, detail="file 不存在或无权访问")
    holder = await _drive_lock_get_holder(file_id)
    if not holder:
        return FileLockResponse(file_id=file_id, locked=False)
    return FileLockResponse(
        file_id=file_id,
        locked=True,
        holder_user_id=int(holder.get("user_id", 0)) or None,
        holder_username=holder.get("username"),
        holder_name=holder.get("name"),
        acquired_at=holder.get("acquired_at"),
        ttl_remaining=int(holder.get("ttl_remaining", 0)),
    )


# ==========================================================
# v2 PR2: 回收站 + 收藏 + 批量操作
# ==========================================================


class TrashListResponse(BaseModel):
    """v2 PR2: 回收站列表响应 (复用 DriveFileListResponse schema)"""
    items: List[DriveFileItem]
    total: int
    page: int
    page_size: int


class FolderStarItem(BaseModel):
    """批次⑩: 收藏列表里的文件夹条目 (表格文件夹行直读)"""
    id: int
    name: str
    owner_id: Optional[int] = None
    owner_name: Optional[str] = None
    owner_avatar: Optional[str] = None
    created_at: Optional[str] = None
    updated_at: Optional[str] = None
    starred_at: Optional[str] = None
    # 子目录下最新文件时间 (前端只取月份显示); 空夹为 None
    latest_file_at: Optional[str] = None
    size_bytes: Optional[int] = None
    is_starred: bool = True


class StarredListResponse(BaseModel):
    """v2 PR2: 收藏列表响应 (批次⑩ 起合并带出收藏的文件夹)"""
    items: List[DriveFileItem]
    total: int
    page: int
    page_size: int
    folders: List[FolderStarItem] = []
    folder_total: int = 0


class ToggleStarResponse(BaseModel):
    """v2 PR2: 收藏切换响应"""
    file_id: int
    is_starred: bool
    starred_at: Optional[str] = None


class BatchIdsRequest(BaseModel):
    """v2 PR2: 通用 batch ids body"""
    file_ids: List[int]


class BatchMoveRequest(BaseModel):
    """v2 PR2: 批量移动请求体"""
    file_ids: List[int]
    target_folder_id: Optional[int] = None  # None = 顶级


class BatchVisibilityRequest(BaseModel):
    """v2 PR2: 批量改可见性请求体"""
    file_ids: List[int]
    visibility: str  # private | team | public


class BatchOperationResponse(BaseModel):
    """v2 PR2: 批量操作统一响应"""
    succeeded_count: int
    skipped_ids: List[int]
    skipped_reasons: Optional[dict] = None  # {file_id: "越权/folder不兼容/不在trash/..."}


class BatchStarRequest(BaseModel):
    """批次①: 批量收藏请求体 (per-user 视角, 幂等)"""
    file_ids: List[int]
    starred: bool = True  # True=收藏, False=取消收藏


class BatchStarResponse(BaseModel):
    """批次①: 批量收藏响应 (updated = 命中目标状态的本人在册 drive 文件数)"""
    updated: int


# ---------- 收藏 ----------

@router.post("/files/batch-star")
async def batch_star_files(
    payload: BatchStarRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """批量收藏/取消收藏 (批次① 收藏个人化, 仅影响**当前成员**自己的收藏夹).

    幂等: 重复收藏同批 id 不产生重复行 (ON CONFLICT DO NOTHING), 取消未收藏的 id
    也不报错; kb 行/不存在 id 静默跳过。返回 {updated: n}。
    """
    svc = DriveService(db)
    updated = await svc.batch_star_files(
        payload.file_ids,
        current_user.id,
        starred=payload.starred,
    )
    return BatchStarResponse(updated=updated)


@router.post("/files/{file_id}/toggle-star", response_model=ToggleStarResponse)
async def toggle_file_star(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """切换文件收藏状态 (批次① 收藏个人化: 仅影响当前成员自己的收藏夹). 404 仅在文件不存在."""
    svc = DriveService(db)
    result = await svc.toggle_star_file(file_id, current_user_id=current_user.id)
    if result is None:
        raise HTTPException(status_code=404, detail="file 不存在")
    f, starred_now, starred_at_now = result
    return ToggleStarResponse(
        file_id=f.id,
        is_starred=starred_now,
        starred_at=str(starred_at_now) if starred_at_now else None,
    )


# ---------- 收藏列表 ----------

@router.get("/starred", response_model=StarredListResponse)
async def list_starred_files(
    page: int = Query(1, ge=1),
    page_size: int = Query(50, ge=1, le=100),
    sort_by: Optional[str] = Query("starred_at", description="starred_at | created_at | updated_at"),
    sort_order: Optional[str] = Query("desc"),
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """列**仅当前成员的个人收藏** (批次① N2 修正 docstring: 134 起 drive_file_stars
    per-user 关系, A 的收藏不出现在 B 的列表; 旧文案"仅 created_by == me"与实现均不符)."""
    svc = DriveService(db)
    items, total = await svc.list_starred(
        current_user_id=current_user.id,
        page=page,
        page_size=page_size,
        sort_by=sort_by or "starred_at",
        sort_order=sort_order or "desc",
    )
    # 列表本身已按本人 star 过滤, is_starred 恒 True, 但仍走同一 attach 路径保证语义单一
    starred_ids = await svc.get_starred_ids([x.id for x in items], current_user.id)
    folder_map = await _folder_name_map(db, items)
    owner_map = await _owner_lookup(db, items)

    # ── 批次⑩: 合并带出本人收藏的文件夹 (表格文件夹行直读; 大小 = path 前缀递归合计) ──
    folder_rows = await db.execute(
        text(
            """
            SELECT f.id, f.name, f.owner_id, m.name AS owner_name, m.avatar AS owner_avatar,
                   f.created_at, f.updated_at, s.starred_at,
                   (SELECT COALESCE(SUM(k.file_size), 0)
                      FROM knowledge k
                      JOIN folders f2 ON k.folder_id = f2.id
                     WHERE f2.path LIKE f.path || '%'
                       AND f2.deleted_at IS NULL
                       AND k.deleted_at IS NULL) AS size_bytes,
                   (SELECT MAX(k.created_at)
                      FROM knowledge k
                      JOIN folders f2 ON k.folder_id = f2.id
                     WHERE f2.path LIKE f.path || '%'
                       AND f2.deleted_at IS NULL
                       AND k.deleted_at IS NULL) AS latest_file_at
              FROM drive_folder_stars s
              JOIN folders f ON f.id = s.folder_id
              JOIN members m ON m.id = f.owner_id
             WHERE s.member_id = :me
             ORDER BY s.starred_at DESC
            """
        ),
        {"me": current_user.id},
    )
    starred_folders = [
        {
            "id": r.id,
            "name": r.name,
            "owner_id": r.owner_id,
            "owner_name": r.owner_name,
            "owner_avatar": r.owner_avatar,
            "created_at": str(r.created_at) if r.created_at else None,
            "updated_at": str(r.updated_at) if r.updated_at else None,
            "starred_at": str(r.starred_at) if r.starred_at else None,
            "latest_file_at": str(r.latest_file_at) if r.latest_file_at else None,
            "size_bytes": int(r.size_bytes or 0),
            "is_starred": True,
        }
        for r in folder_rows
    ]

    return StarredListResponse(
        items=[_to_item(x, starred_ids=starred_ids, folder_map=folder_map, owner_lookup=owner_map) for x in items],
        total=total,
        page=page,
        page_size=page_size,
        folders=starred_folders,
        folder_total=len(starred_folders),
    )


# ---------- 回收站 ----------

@router.get("/trash", response_model=TrashListResponse)
async def list_trash_files(
    page: int = Query(1, ge=1),
    page_size: int = Query(50, ge=1, le=100),
    sort_by: Optional[str] = Query("deleted_at"),
    sort_order: Optional[str] = Query("desc"),
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """列回收站文件 (软删的 drive 文件, 仅 owner)."""
    svc = DriveService(db)
    items, total = await svc.list_trash(
        current_user_id=current_user.id,
        page=page,
        page_size=page_size,
        sort_by=sort_by or "deleted_at",
        sort_order=sort_order or "desc",
    )
    owner_map = await _owner_lookup(db, items)
    return TrashListResponse(
        items=[_to_item(x, owner_lookup=owner_map) for x in items],
        total=total,
        page=page,
        page_size=page_size,
    )


@router.post("/trash/permanent-delete", response_model=BatchOperationResponse)
async def permanent_delete_files(
    payload: BatchIdsRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """批量物理删除回收站中的文件 (不可逆; 2026-09-05 角色扁平化: 任何成员可操作)."""
    svc = DriveService(db)
    deleted, skipped = await svc.permanent_delete_batch(
        payload.file_ids,
        current_user_id=current_user.id,
        is_admin=True,
    )
    return BatchOperationResponse(
        succeeded_count=deleted,
        skipped_ids=skipped,
        skipped_reasons={fid: "不在回收站/不存在/非 owner" for fid in skipped},
    )


# ---------- 批量操作 ----------

@router.post("/files/batch-soft-delete", response_model=BatchOperationResponse)
async def batch_soft_delete_files(
    payload: BatchIdsRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """批量软删 (进入回收站)."""
    svc = DriveService(db)
    deleted, skipped = await svc.batch_soft_delete(
        payload.file_ids, current_user_id=current_user.id,
    )
    return BatchOperationResponse(
        succeeded_count=deleted,
        skipped_ids=skipped,
        skipped_reasons={fid: "不存在/非 owner" for fid in skipped},
    )


@router.post("/files/batch-restore", response_model=BatchOperationResponse)
async def batch_restore_files(
    payload: BatchIdsRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """批量从回收站恢复 (2026-09-05 角色扁平化: 任何成员可操作)."""
    svc = DriveService(db)
    restored, skipped = await svc.batch_restore(
        payload.file_ids,
        current_user_id=current_user.id,
        is_admin=True,
    )
    return BatchOperationResponse(
        succeeded_count=restored,
        skipped_ids=skipped,
        skipped_reasons={fid: "不在回收站/不存在/非 owner" for fid in skipped},
    )


@router.post("/files/batch-move", response_model=BatchOperationResponse)
async def batch_move_files(
    payload: BatchMoveRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """批量移动到 folder (target_folder_id=None = 顶级)."""
    svc = DriveService(db)
    try:
        moved, skipped = await svc.batch_move(
            payload.file_ids,
            target_folder_id=payload.target_folder_id,
            current_user_id=current_user.id,
        )
    except DriveServiceError as e:
        # W1 T1 migration: DriveServiceError → AppException envelope
        raise_app_error(e.status_code, ERR_FILE_FORBIDDEN, str(e))
    return BatchOperationResponse(
        succeeded_count=moved,
        skipped_ids=skipped,
        skipped_reasons={fid: "folder 上限/不存在/非 owner" for fid in skipped},
    )


@router.post("/files/batch-update-visibility", response_model=BatchOperationResponse)
async def batch_update_files_visibility(
    payload: BatchVisibilityRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """批量改可见性 (folder 上限校验)."""
    svc = DriveService(db)
    try:
        updated, skipped = await svc.batch_update_visibility(
            payload.file_ids,
            new_visibility=payload.visibility,
            current_user_id=current_user.id,
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=str(e))
    return BatchOperationResponse(
        succeeded_count=updated,
        skipped_ids=skipped,
        skipped_reasons={fid: "folder 上限/不存在/非 owner" for fid in skipped},
    )


# === 容量统计 ===

@router.get("/storage-stats", response_model=StorageStatsResponse)
async def get_storage_stats(
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """当前用户的 drive 存储统计 (file_count + by_visibility)"""
    svc = DriveService(db)
    stats = await svc.storage_stats(current_user_id=current_user.id)
    return StorageStatsResponse(**stats)


# ==========================================================================
# 下载 (PR2.6)
# ==========================================================================
import zipfile
import re
from urllib.parse import quote

from fastapi.responses import StreamingResponse
from starlette.requests import Request


def build_content_disposition(disposition: str, filename: str) -> str:
    """构建 Content-Disposition 头 (RFC 5987 + RFC 6266).

    历史教训 (2026-07-02):
      旧实现用 `filename="{filename}"; filename*=UTF-8''{encoded}`
      双 attribute 看似稳妥, 实则 `filename=` 部分走 latin-1 codec,
      一旦 filename 含中文 (如 "组会ppt/冯懿鑫/2025.7.2 研一 冯懿鑫.pptx"),
      Starlette/FastAPI 调用 latin-1 encode → UnicodeEncodeError → 500.

    修复: 仅输出 `filename*=UTF-8''<encoded>` (RFC 5987 标准化形式),
    移除非 ASCII safe 的 `filename="..."` 旧 attribute.
    现代浏览器 (Chrome / Firefox / Safari / Edge) 全部支持 filename*,
    老 IE (≤ IE9) 不支持但本项目目标用户无 IE.

    Args:
        disposition: 'attachment' 或 'inline'
        filename: 原始文件名 (可含中文)

    Returns:
        e.g. "inline; filename*=UTF-8''%E7%BB%84%E4%BC%9Appt.pptx"
    """
    encoded = quote(filename, safe='')
    return f"{disposition}; filename*=UTF-8''{encoded}"


def _check_download_visibility(file_knowledge, current_user_id: int) -> None:
    """下载/预览前校验 — 2026-09 单一团队空间: pass-through (私有下载门禁已退役)。

    函数与调用点保留 (防签名破坏); 迁移 133 回填后 drive 行无 private, 即便有脏数据
    也不再拦截。
    """
    return None


@router.get("/files/{file_id}/download")
async def download_drive_file(
    file_id: int,
    request: Request,
    disposition: str = Query("attachment", pattern="^(attachment|inline)$"),
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """单文件下载 (支持 Range 断点续传 + 中文文件名 URL 编码)"""
    svc = DriveService(db)
    f = await svc.get_file(file_id, current_user_id=current_user.id)
    if f is None:
        raise HTTPException(status_code=404, detail="file 不存在或无权访问")
    _check_download_visibility(f, current_user.id)

    if not f.file_path:
        raise HTTPException(status_code=404, detail="file 无 MinIO 对象")

    # 文件元信息
    # PR4.6 修复: f.file_type 是 ".txt" 字面量 (DB 存的是扩展名), 直接返会触发 nosniff 阻断
    # 用 mimetypes 模块从扩展名推断真实 MIME (image/jpeg, video/mp4, application/pdf 等)
    import mimetypes
    if f.file_type and f.file_type.startswith("."):
        # f.file_type 是 ".ext" 形式 (我们 KB 上传时存的就是这个)
        guessed, _ = mimetypes.guess_type(f"a{f.file_type}")
        content_type = guessed or "application/octet-stream"
    else:
        content_type = f.file_type or "application/octet-stream"

    if content_type and not content_type.startswith("video/") and not content_type.startswith("image/") and disposition == "inline":
        # text 类保持原 mime, 其他 inline 也按 octet-stream
        pass
    filename = f.file_name or f.title or f"file_{f.id}"

    # PR2.7: 原子 +1 下载计数
    new_count = await svc.increment_download_count(file_id)

    # Range 头解析
    range_header = request.headers.get("Range")
    if range_header:
        m = re.match(r"bytes=(\d+)-(\d*)", range_header)
        if m:
            start = int(m.group(1))
            end_str = m.group(2)
            # 读 total size
            file_size = await _get_object_size(f.file_path)
            if file_size is None:
                # 没拿到 size, 走完整下载
                start = None
            else:
                if end_str:
                    end = min(int(end_str), file_size - 1)
                else:
                    end = file_size - 1
                length = end - start + 1

                async def _range_stream():
                    chunk = await _download_range(f.file_path, start, length)
                    yield chunk

                headers = {
                    "Content-Range": f"bytes {start}-{end}/{file_size}",
                    "Accept-Ranges": "bytes",
                    "Content-Length": str(length),
                    "Content-Disposition": build_content_disposition(disposition, filename),
                }
                return StreamingResponse(
                    _range_stream(),
                    status_code=206,
                    media_type=content_type,
                    headers=headers,
                )

    # 完整下载
    file_size = await _get_object_size(f.file_path)
    headers = {
        "Content-Disposition": build_content_disposition(disposition, filename),
    }
    if file_size is not None:
        headers["Content-Length"] = str(file_size)
    headers["Accept-Ranges"] = "bytes"

    async def _full_stream():
        data = await file_service.download_file(f.file_path)
        yield data

    return StreamingResponse(
        _full_stream(),
        media_type=content_type,
        headers=headers,
    )


async def _get_object_size(object_name: str) -> Optional[int]:
    """读 MinIO 对象 size (用 stat_object 同步)"""
    import asyncio
    try:
        def _sync_stat():
            return file_service.client.stat_object(file_service.bucket, object_name)
        obj = await asyncio.to_thread(_sync_stat)
        return obj.size
    except Exception:
        return None


async def _download_range(object_name: str, start: int, length: int) -> bytes:
    """下载 MinIO 对象指定字节范围 (minio get_object 支持 offset+length)"""
    import asyncio
    def _sync_range():
        response = file_service.client.get_object(
            file_service.bucket, object_name, offset=start, length=length,
        )
        try:
            return response.read()
        finally:
            response.close()
            response.release_conn()
    return await asyncio.to_thread(_sync_range)


# === 批量 ZIP 下载 ===

class BatchDownloadRequest(BaseModel):
    ids: Optional[List[int]] = Field(None, description="文件 id 列表")
    folder_id: Optional[int] = Field(None, description="递归下载整个 folder")


@router.post("/files/batch-download")
async def batch_download_drive_files(
    payload: BatchDownloadRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """批量 ZIP 下载 (流式生成, 不落盘)

    body: {"ids": [1,2,3]} OR {"folder_id": 4}
    无权限文件跳过, ZIP 根目录生成 _skipped.txt
    """
    if not payload.ids and not payload.folder_id:
        # W1 T1 migration: 400 → AppException envelope (BATCH_PARAM_MISSING)
        raise_app_error(400, ERR_BATCH_PARAM_MISSING, "ids 或 folder_id 必填其一")

    svc = DriveService(db)

    # 1) 收集 file 列表
    file_records = []
    skipped = []
    if payload.ids:
        for fid in payload.ids:
            f = await svc.get_file(fid, current_user_id=current_user.id)
            if f is None:
                skipped.append(f"id={fid} 无权访问")
                continue
            _check_download_visibility(f, current_user.id)
            if not f.file_path:
                skipped.append(f"id={fid} 无 MinIO 对象")
                continue
            file_records.append(f)
    elif payload.folder_id:
        # 递归收集 folder 下的所有文件 (含子 folder)
        file_records, skipped = await _collect_folder_files(
            db, svc, payload.folder_id, current_user.id, skipped
        )

    if not file_records and not skipped:
        # W1 T1 migration: 404 → AppException envelope
        raise_app_error(404, ERR_FILE_NOT_FOUND, "无可下载文件")

    # 2) 流式生成 ZIP
    timestamp = _dt_now.now().strftime("%Y%m%d_%H%M%S")
    zip_filename = f"drive_{current_user.username}_{timestamp}.zip"

    async def _zip_stream():
        # BytesIO 缓冲区 zip stream
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for f in file_records:
                try:
                    data = await file_service.download_file(f.file_path)
                    # ZIP 内路径: file_name (避免重复)
                    arcname = f.file_name or f"file_{f.id}"
                    # 防路径冲突 (同名)
                    existing = [n for n in zf.namelist() if n == arcname]
                    if existing:
                        arcname = f"{f.id}_{arcname}"
                    zf.writestr(arcname, data)
                except Exception as e:
                    logger.warning(f"批量下载跳过 file_id={f.id}: {e}")
                    continue
            # skipped list 写 _skipped.txt
            if skipped:
                skipped_content = "\n".join(skipped)
                zf.writestr("_skipped.txt", skipped_content)
        # yield 一次完整 buffer (PR3 优化: chunked write)
        buffer.seek(0)
        yield buffer.read()

    return StreamingResponse(
        _zip_stream(),
        media_type="application/zip",
        headers={
            "Content-Disposition": build_content_disposition("attachment", zip_filename),
        },
    )


async def _collect_folder_files(
    db: AsyncSession,
    drive_svc: DriveService,
    folder_id: int,
    current_user_id: int,
    skipped: list,
) -> tuple:
    """递归收集 folder 下所有 file + 跨子 folder, 越权文件入 skipped"""
    from app.services.folder_service import FolderService
    folder_svc = FolderService(db)

    # 1) 校验 folder 存在 (2026-09 单一团队空间: 无 private/owner 越权跳过)
    folder = await folder_svc.get_folder(folder_id)
    if folder is None:
        return [], skipped + [f"folder_id={folder_id} 不存在"]

    # 2) 列当前 folder 的文件
    files, _ = await drive_svc.list_files(
        current_user_id=current_user_id,
        folder_id=folder_id,
        storage_mode="drive",
        page=1,
        page_size=1000,  # 单 folder 上限 1000 个文件
    )
    # 过滤已软删 + 真实可见
    file_records = [f for f in files if f.deleted_at is None and f.file_path]

    # 3) 递归子 folder (2026-09 单一团队空间: 子 folder 不再按 private/owner 跳过)
    children = await folder_svc.list_children(folder_id=folder_id, include_deleted=False)
    for child in children:
        sub_files, skipped = await _collect_folder_files(
            db, drive_svc, child.id, current_user_id, skipped,
        )
        file_records.extend(sub_files)

    return file_records, skipped


# === Stats 补 time import ===
from datetime import datetime as _dt_now


# ==========================================================================
# PR2.7 分享链接 + 公开下载
# ==========================================================================

class ShareLinkRequest(BaseModel):
    """v2 PR1: 新增 expires_hours (细粒度) + password (提取码) 字段

    expires_in_days 保留旧字段向后兼容 (优先级低于 expires_hours).
    password 必填 4-8 位数字, 后端存 SHA256 hash (不存明文).
    """
    expires_in_days: Optional[int] = None  # 旧 API, 1-365
    expires_hours: Optional[int] = None    # 新 API, 1-8760; 0=默认 7 天; -1=永久
    password: Optional[str] = None          # 4-8 位数字, None=无密码


class ShareLinkResponse(BaseModel):
    file_id: int
    token: str
    share_url: str
    expires_at: Optional[str] = None
    password_required: bool = False         # v2 PR1 新增: 是否需要提取码


@router.post("/files/{file_id}/share-link", response_model=ShareLinkResponse)
async def create_share_link(
    file_id: int,
    payload: ShareLinkRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """生成 drive 文件公开分享链接 (2026-09 单一团队空间: 任何成员可分享)

    返回 token + 完整 share_url + expires_at + password_required
    """
    svc = DriveService(db)
    try:
        f = await svc.create_share_link(
            file_id,
            current_user_id=current_user.id,
            expires_in_days=payload.expires_in_days,
            expires_hours=payload.expires_hours,
            password=payload.password,
        )
    except DriveServiceError as e:
        # W1 T1 migration: DriveServiceError → AppException envelope
        raise_app_error(e.status_code, ERR_FILE_FORBIDDEN, str(e), file_id=file_id)
    if f is None:
        raise_app_error(404, ERR_SHARE_LINK_NOT_FOUND, "file 不存在", file_id=file_id)

    # share_url 用 settings.PUBLIC_BASE_URL (前端拼) - 这里用相对路径
    share_url = f"/drive/share/{f.share_token}"
    return ShareLinkResponse(
        file_id=f.id,
        token=f.share_token,
        share_url=share_url,
        expires_at=str(f.share_expires_at) if f.share_expires_at else None,
        password_required=bool(f.share_password),
    )


@router.delete("/files/{file_id}/share-link", status_code=204)
async def revoke_share_link(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """撤销 share link (清 token)"""
    svc = DriveService(db)
    ok = await svc.revoke_share_link(file_id, current_user_id=current_user.id)
    if not ok:
        raise_app_error(404, ERR_SHARE_LINK_NOT_FOUND, "file 不存在", file_id=file_id)
    return


# === v2 PR1 公开分享 GET 端点 (含密码验证) ===



# === 批次⑩.17 自研 PPT 第三栏预览: python-pptx 解析为结构化 JSON (2026-09-06) ===
# .pptx = zip + OOXML — python-pptx 已在容器内 (1.0.2), 解析一次缓存 JSON,
# 第三栏渲染器 (DriveDetailRail) 按 EMU 比例绝对定位还原 文本框/图片/表格。
_PPTX_STRUCTURE_CACHE: dict = {}  # (file_id, updated_at) -> payload (保留最近 16 份)


@router.get("/files/{file_id}/pptx-structure")
async def get_pptx_structure(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """自研 PPT 预览数据源: python-pptx 抽取 幻灯片结构 JSON (无 JWT 不可访问).

    返回: {total, slide_w_emu, slide_h_emu, slides: [{index, bg, shapes: [...]}]}
    shape.kind: text (runs 带 字号pt/颜色/粗体) | image (dataURL, 单张≤1.5MB) |
                table (rows) | chart (占位, 二期 ECharts 重绘) | shape (纯色块)
    坐标/尺寸均为 0-1 比例 (相对幻灯片宽高), 前端按舞台宽度缩放。
    """
    import base64
    from io import BytesIO
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    svc = DriveService(db)
    f = await svc.get_file(file_id, current_user_id=current_user.id)
    if f is None:
        raise HTTPException(status_code=404, detail="file 不存在或无权访问")
    if not (f.file_name or "").lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="仅支持 .pptx 文件")
    if not f.file_path:
        raise HTTPException(status_code=404, detail="file 无 MinIO 对象")

    cache_key = (file_id, str(f.updated_at))
    if cache_key in _PPTX_STRUCTURE_CACHE:
        return _PPTX_STRUCTURE_CACHE[cache_key]

    raw = await file_service.download_file(f.file_path)
    try:
        prs = Presentation(BytesIO(raw))
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"pptx 解析失败: {e}")

    slide_w, slide_h = prs.slide_width or 9144000, prs.slide_height or 6858000

    # ── 主题色板 (theme1.xml clrScheme) ──
    import re as _re
    import zipfile
    theme = {}
    try:
        with zipfile.ZipFile(BytesIO(raw)) as zf:
            tname = next((n for n in zf.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")), None)
            if tname:
                tseg = zf.read(tname).decode("utf-8", "ignore")
                mm = _re.search(r"<a:clrScheme.*?</a:clrScheme>", tseg, _re.S)
                if mm:
                    for km in _re.finditer(r"<a:(\w+)>(.*?)</a:\1>", mm.group(0), _re.S):
                        cm = _re.search(r'val="([0-9A-Fa-f]{6})"', km.group(2)) or _re.search(r'lastClr="([0-9A-Fa-f]{6})"', km.group(2))
                        if cm:
                            theme[km.group(1).lower()] = cm.group(1).lower()
    except Exception:
        theme = {}

    def _theme_hex(enum_name):
        k = str(enum_name).split(" ")[0].lower()
        k = {"dark_1": "dk1", "text_1": "dk1", "light_1": "lt1",
             "background_1": "lt1", "text_2": "dk2", "background_2": "lt2"}.get(k, k.replace("_", ""))
        return theme.get(k)

    def _scheme_hex(val):
        k = str(val).lower()
        k = {"tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2"}.get(k, k)
        return theme.get(k)

    # ── 母版/版式 默认字号与颜色 (错版根因: 占位符文字大量无显式样式) ──
    from pptx.oxml.ns import qn as _qn
    from PIL import Image as PILImage

    def _lst1(txbody_el):
        """txBody 的 lvl1pPr defRPr 元素 (可能为 None)"""
        try:
            ls = txbody_el.find(_qn("a:lstStyle"))
            if ls is not None:
                return ls.find(_qn("a:lvl1pPr"))
        except Exception:
            pass
        return None

    def _dr_sz(dr):
        try:
            v = dr.get("sz")
            return int(v) / 100.0 if v else None
        except Exception:
            return None

    def _dr_color(dr):
        try:
            sf = dr.find(_qn("a:solidFill"))
            if sf is None:
                return None
            srgb = sf.find(_qn("a:srgbClr"))
            if srgb is not None:
                return srgb.get("val")
            sch = sf.find(_qn("a:schemeClr"))
            if sch is not None:
                return _scheme_hex(sch.get("val"))
        except Exception:
            pass
        return None

    def _lvl_defrpr_para(para):
        try:
            if para._pPr is not None:
                return para._pPr.find(_qn("a:defRPr"))
        except Exception:
            pass
        return None

    def _lvl_defrpr(lst_el, lvl):
        try:
            if lst_el is None:
                return None
            name = "a:lvl%dpPr" % min(max(lvl, 1), 9)
            l1 = lst_el.find(_qn(name))
            if l1 is not None:
                return l1.find(_qn("a:defRPr"))
        except Exception:
            pass
        return None

    def _tf_defaults(txbody_el):
        dr = _lst1(txbody_el)
        if dr is None:
            return (None, None)
        return (_dr_sz(dr), _dr_color(dr))

    def _ph_defaults(ph, palette):
        """版式占位符默认 (字号pt, 颜色): lstStyle lvl1 → 母版同型占位符 lstStyle → 母版 txStyles"""
        try:
            sz, color = _tf_defaults(ph.text_frame._txBody)
            if sz is None or color is None:
                ptype = str(ph.placeholder_format.type)
                is_title = "TITLE" in ptype or "CENTER" in ptype
                master = ph.part.slide_layout.slide_master if hasattr(ph.part, "slide_layout") else None
                if master is None:
                    try:
                        master = ph.slide_layout.slide_master
                    except Exception:
                        master = None
                if master is not None:
                    if sz is None:
                        for mph in master.placeholders:
                            mpt = str(mph.placeholder_format.type)
                            if ("TITLE" in mpt) == is_title and "TITLE" in mpt:
                                s2, c2 = _tf_defaults(mph.text_frame._txBody)
                                sz = sz or s2
                                color = color or c2
                                break
                    if (sz is None or color is None) and master.element is not None:
                        txs = master.element.find(_qn("p:txStyles"))
                        if txs is not None:
                            style_el = txs.find(_qn("p:titleStyle" if is_title else "p:bodyStyle"))
                            if style_el is None and not is_title:
                                style_el = txs.find(_qn("p:otherStyle"))
                            if style_el is not None:
                                l1 = style_el.find(_qn("a:lvl1pPr"))
                                if l1 is not None:
                                    dr = l1.find(_qn("a:defRPr"))
                                    if dr is not None:
                                        sz = sz or _dr_sz(dr)
                                        color = color or _dr_color(dr)
            if color is None:
                color = palette.get("dk1")
            return (sz, color)
        except Exception:
            return (None, None)

    def _fill_hex(sh, palette):
        try:
            fill = sh.fill
            if fill.type is None or int(fill.type) != 1:
                return None
            fc = fill.fore_color
            try:
                if fc.type is not None and "RGB" in str(fc.type):
                    return str(fc.rgb)
            except Exception:
                pass
            try:
                return _theme_hex(str(fc.theme_color))
            except Exception:
                return None
        except Exception:
            return None

    def _layout_ph_map(slide, palette):
        m = {}
        try:
            for ph in slide.slide_layout.placeholders:
                m[ph.placeholder_format.idx] = (_ph_defaults(ph, palette), ph)
        except Exception:
            pass
        return m

    def _inherit_pos(sh, lay_map):
        try:
            if sh.left is not None or not getattr(sh, "is_placeholder", False):
                return None
            entry = lay_map.get(sh.placeholder_format.idx)
            ph = entry[1] if entry else None
            if ph is None:
                ptype = sh.placeholder_format.type
                match = next((e for k, e in lay_map.items()
                              if e[1].placeholder_format.type == ptype), None)
                ph = match[1] if match else None
            if ph is not None and ph.left is not None:
                return (ph.left, ph.top, ph.width or 0, ph.height or 0)
        except Exception:
            pass
        return None

    def _run_size_pt(sh, lay_map, para, shape_default):
        if para is not None and para.font.size is not None:
            return para.font.size.pt
        if para is not None and para._pPr is not None:
            ps = _dr_sz(_lvl_defrpr_para(para))
            if ps:
                return ps
        lvl = (para.level if para is not None and para.level else 0) + 1
        try:
            lst_el = sh.text_frame._txBody.find(_qn("a:lstStyle"))
            ps = _dr_sz(_lvl_defrpr(lst_el, lvl))
            if ps:
                return ps
        except Exception:
            pass
        try:
            if getattr(sh, "is_placeholder", False):
                entry = lay_map.get(sh.placeholder_format.idx)
                if entry:
                    ph, lst_el = entry[1], entry[1].text_frame._txBody
                    ps = _dr_sz(_lvl_defrpr(lst_el, lvl))
                    if ps:
                        return ps
        except Exception:
            pass
        try:
            master = slide.slide_layout.slide_master
            txs = master.element.find(_qn("p:txStyles"))
            if txs is not None:
                is_title = "TITLE" in str(sh.placeholder_format.type) if getattr(sh, "is_placeholder", False) else False
                style_el = txs.find(_qn("p:titleStyle" if is_title else "p:bodyStyle"))
                if style_el is not None:
                    ps = _dr_sz(_lvl_defrpr(style_el, 1 if is_title else lvl))
                    if ps:
                        return ps
        except Exception:
            pass
        if shape_default:
            return shape_default
        return 14

    def _para_defrpr(para):
        try:
            if para is not None and para._pPr is not None:
                return para._pPr.find(_qn("a:defRPr"))
        except Exception:
            pass
        return None

    def _run_color(sh, lay_map, para, shape_default_color):
        if para is not None and para.font.color is not None and para.font.color.type is not None:
            try:
                if "RGB" in str(para.font.color.type):
                    return str(para.font.color.rgb)
                return _theme_hex(str(para.font.color.theme_color))
            except Exception:
                pass
        lvl = (para.level if para is not None and para.level else 0) + 1
        dr = _para_defrpr(para)
        if dr is not None:
            c = _dr_color(dr)
            if c:
                return c
        try:
            lst_el = sh.text_frame._txBody.find(_qn("a:lstStyle"))
            dr = _lvl_defrpr(lst_el, lvl)
            if dr is not None:
                c = _dr_color(dr)
                if c:
                    return c
        except Exception:
            pass
        try:
            if getattr(sh, "is_placeholder", False):
                entry = lay_map.get(sh.placeholder_format.idx)
                if entry:
                    dr = _lvl_defrpr(entry[1].text_frame._txBody, lvl)
                    if dr is not None:
                        c = _dr_color(dr)
                        if c:
                            return c
        except Exception:
            pass
        if shape_default_color:
            return shape_default_color
        try:
            master = slide.slide_layout.slide_master
            txs = master.element.find(_qn("p:txStyles"))
            if txs is not None:
                is_title = "TITLE" in str(sh.placeholder_format.type) if getattr(sh, "is_placeholder", False) else False
                style_el = txs.find(_qn("p:titleStyle" if is_title else "p:bodyStyle"))
                if style_el is not None:
                    dr = _lvl_defrpr(style_el, 1 if is_title else lvl)
                    if dr is not None:
                        c = _dr_color(dr)
                        if c:
                            return c
        except Exception:
            pass
        return theme.get("dk1")

    def _downscale_blob(blob, content_type):
        """大图降采样 → JPEG (修「大图跳过」缺内容错版)"""
        try:
            img = PILImage.open(BytesIO(blob))
            img.thumbnail((1280, 1280))
            buf = BytesIO()
            if img.mode in ("RGBA", "P", "LA"):
                img = img.convert("RGB")
            img.save(buf, "JPEG", quality=82)
            out = buf.getvalue()
            if len(out) < len(blob):
                return "data:image/jpeg;base64," + base64.b64encode(out).decode()
        except Exception:
            pass
        return None

    slides = []

    def collect(shapes, fx, fy, fsx, fsy, lay_map, palette, depth):
        out = []
        if depth > 4:
            return out
        for sh in shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        xfrm = sh._element.grpSpPr.xfrm
                        ox, oy = (xfrm.off.x or 0), (xfrm.off.y or 0)
                        ecx, ecy = (xfrm.ext.cx or 1), (xfrm.ext.cy or 1)
                        cox, coy = (xfrm.chOff.x or 0), (xfrm.chOff.y or 0)
                        ccx = (xfrm.chExt.cx or ecx) or 1
                        ccy = (xfrm.chExt.cy or ecy) or 1
                        nfx = fx + fsx * (ox - cox * (ecx / ccx))
                        nfy = fy + fsy * (oy - coy * (ecy / ccy))
                        out += collect(sh.shapes, nfx, nfy, fsx * (ecx / ccx), fsy * (ecy / ccy),
                                       lay_map, palette, depth + 1)
                    except Exception:
                        pass
                    continue

                inherit = _inherit_pos(sh, lay_map)
                left = sh.left if sh.left is not None else (inherit[0] if inherit else 0)
                top = sh.top if sh.top is not None else (inherit[1] if inherit else 0)
                width = sh.width if sh.width is not None else (inherit[2] if inherit else 0)
                height = sh.height if sh.height is not None else (inherit[3] if inherit else 0)

                x = (fx + left * fsx) / slide_w
                y = (fy + top * fsy) / slide_h
                w = (width * fsx) / slide_w
                h = (height * fsy) / slide_h
                if w <= 0 or h <= 0:
                    continue
                entry = {"x": round(x, 4), "y": round(y, 4), "w": round(w, 4), "h": round(h, 4)}

                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob = sh.image.blob
                    src = None
                    if len(blob) <= 1536 * 1024:
                        src = "data:" + (sh.image.content_type or "image/png") + ";base64," + base64.b64encode(blob).decode()
                    else:
                        src = _downscale_blob(blob, sh.image.content_type)
                    if src:
                        entry.update(kind="image", src=src)
                    else:
                        entry.update(kind="image", skip="large")
                elif getattr(sh, "has_table", False) and sh.has_table:
                    tbl = sh.table
                    rows = [[cell.text for cell in row.cells] for row in tbl.rows]
                    entry.update(kind="table", rows=rows, first_row=bool(tbl.first_row))
                elif getattr(sh, "has_chart", False) and sh.has_chart:
                    ctype = "bar"
                    try:
                        ctype = str(sh.chart.chart_type).split(" ")[0].lower()
                    except Exception:
                        pass
                    cats, series = [], []
                    try:
                        cats = [str(c) for c in sh.chart.plots[0].categories]
                    except Exception:
                        pass
                    try:
                        for sdata in sh.chart.plots[0].series:
                            series.append({
                                "name": str(sdata.name or ""),
                                "values": [round(v, 2) if isinstance(v, (int, float)) else None for v in sdata.values],
                            })
                    except Exception:
                        pass
                    entry.update(kind="chart-data", chart_type=ctype, categories=cats, series=series)
                elif sh.has_text_frame and sh.text_frame.text.strip():
                    # 形状级默认字号/颜色: 自身 lstStyle → 版式占位符 → 母版
                    shape_sz, shape_color = None, None
                    try:
                        shape_sz, shape_color = _tf_defaults(sh.text_frame._txBody)
                    except Exception:
                        pass
                    try:
                        if getattr(sh, "is_placeholder", False):
                            entry0 = lay_map.get(sh.placeholder_format.idx)
                            if entry0:
                                shape_sz = shape_sz or entry0[0][0]
                                shape_color = shape_color or entry0[0][1]
                    except Exception:
                        pass
                    if shape_color is None:
                        shape_color = palette.get("dk1") or "1F2A26"

                    paras = []
                    for para in sh.text_frame.paragraphs:
                        runs = []
                        for r in para.runs:
                            color = None
                            try:
                                if r.font.color is not None and r.font.color.type is not None:
                                    if "RGB" in str(r.font.color.type):
                                        color = str(r.font.color.rgb)
                                    else:
                                        color = _theme_hex(str(r.font.color.theme_color))
                            except Exception:
                                color = None
                            sz = r.font.size.pt if r.font.size is not None else _run_size_pt(sh, lay_map, para, shape_sz)
                            color = color or _run_color(sh, lay_map, para, shape_color)
                            runs.append({"t": r.text, "sz": sz, "b": bool(r.font.bold), "c": color})
                        if runs:
                            ls = para.line_spacing
                            paras.append({
                                "runs": runs,
                                "align": str(para.alignment).split(" ")[0].split(":")[-1].lower()
                                if para.alignment is not None else None,
                                "ls": round(float(ls), 2) if isinstance(ls, float) else None,
                            })
                    entry.update(kind="text", paras=paras)
                else:
                    entry.update(kind="shape", color=_fill_hex(sh, palette))
                out.append(entry)
            except Exception:
                continue
        return out

    for idx, slide in enumerate(prs.slides, 1):
        bg = None
        try:
            fill = slide.background.fill
            if fill.type is not None and int(fill.type) == 1:
                bg = str(fill.fore_color.rgb)
        except Exception:
            bg = None
        lay_map = _layout_ph_map(slide, theme)
        shapes = collect(slide.shapes, 0, 0, 1, 1, lay_map, theme, 0)
        slides.append({"index": idx, "bg": bg, "shapes": shapes})

    payload = {
        "file_name": f.file_name,
        "total": len(slides),
        "slide_w_emu": slide_w,
        "slide_h_emu": slide_h,
        "slides": slides,
    }
    if len(_PPTX_STRUCTURE_CACHE) > 16:
        _PPTX_STRUCTURE_CACHE.pop(next(iter(_PPTX_STRUCTURE_CACHE)))
    _PPTX_STRUCTURE_CACHE[cache_key] = payload
    return payload


class PublicShareInfoResponse(BaseModel):
    """公开分享链接的元信息 (验证密码后才返回下载链接)."""
    file_name: str
    file_size: Optional[int] = None
    expires_at: Optional[str] = None
    password_required: bool = False
    verify_token: Optional[str] = None  # 验证密码后由后端签发的短期 token, 传给下载端点


class PublicShareDownloadRequest(BaseModel):
    """下载请求 body. password 留空时直接尝试下载 (公开分享无密码模式)."""
    password: Optional[str] = None


# === share_router endpoints 注册放在 share_router 定义之后 (顺序依赖) ===
# 见下方 share_router 段


# ==========================================================================
# v2 PR1 visibility edit endpoint (桌面 stub 修复)
# ==========================================================================

class UpdateVisibilityRequest(BaseModel):
    """修改 drive 文件可见性请求体."""
    visibility: str  # private | team | public


class UpdateVisibilityResponse(BaseModel):
    file_id: int
    visibility: str
    folder_id: Optional[int] = None


@router.put("/files/{file_id}/visibility", response_model=UpdateVisibilityResponse)
async def update_file_visibility(
    file_id: int,
    payload: UpdateVisibilityRequest,
    db: AsyncSession = Depends(get_db),
    current_user: Member = Depends(get_current_user),
):
    """PUT /api/v1/drive/files/{file_id}/visibility

    修改 drive 文件可见性 (owner only, 校验文件夹硬上限).
    桌面 DriveView 的 handleFileUpdateVisibility stub 修复由此端点支撑.
    """
    svc = DriveService(db)
    try:
        f = await svc.update_visibility(
            file_id,
            current_user_id=current_user.id,
            new_visibility=payload.visibility,
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=str(e))
    if f is None:
        raise HTTPException(status_code=404, detail="file 不存在")
    return UpdateVisibilityResponse(
        file_id=f.id,
        visibility=f.visibility,
        folder_id=f.folder_id,
    )


# === 公开 share 端点 (无需 JWT, token 验证) ===

# 注: 放在前缀 /drive/share 路径下, router prefix = /drive
# 实际 URL: GET /api/v1/drive/share/{token}
# 为简化, 把 share 端点放到 root 上
# 但 FastAPI router prefix 是 /drive, 加 /share/{token} 即可
# 这里用独立 APIRouter 避免与 /files/{file_id} 冲突
share_router = APIRouter(prefix="/drive/share", tags=["网盘公开分享"])


@share_router.get("/{token}/info", response_model=PublicShareInfoResponse)
async def public_share_info(
    token: str,
    db: AsyncSession = Depends(get_db),
):
    """GET /api/v1/drive/share/{token}/info — 查看分享链接元信息 + 是否需要密码.

    注意: 此端点不返回密码本身, 仅返回 password_required flag + 文件元信息.
    """
    svc = DriveService(db)
    f = await svc.get_by_share_token(token)
    if f is None:
        raise HTTPException(status_code=404, detail="分享链接不存在或已过期")
    return PublicShareInfoResponse(
        file_name=f.file_name or f.title or f"file_{f.id}",
        file_size=None,  # 暂不返具体 size (避免被白嫖)
        expires_at=str(f.share_expires_at) if f.share_expires_at else None,
        password_required=bool(f.share_password),
    )


@share_router.post("/{token}/verify-password")
async def public_share_verify_password(
    token: str,
    request: PublicShareDownloadRequest,
    db: AsyncSession = Depends(get_db),
):
    """POST /api/v1/drive/share/{token}/verify-password

    验证提取码. 失败返 403 (不放任何提示防枚举). 成功返 verify_token (短期有效).
    """
    svc = DriveService(db)
    f = await svc.verify_share_access(token, password=request.password)
    if f is None:
        # 整体返 403 不区分"不存在/过期/密码错" 防 token 枚举
        raise HTTPException(status_code=403, detail="分享链接已过期或密码错误")
    # v2 PR1 暂用文件 id 作为 verify_token (无 JWT 时跳转用)
    return {"verified": True, "file_id": f.id}


@share_router.get("/{token}")
async def public_download_by_token(
    token: str,
    request: Request,
    password: Optional[str] = Query(None, description="提取码 (分享链接有密码时必填)"),
    db: AsyncSession = Depends(get_db),
):
    """公开分享下载 (无 JWT, 校验 token + 可选提取码)

    GET /api/v1/drive/share/{token}?password=1234 -> 流式下载
    GET /api/v1/drive/share/{token}              -> 无密码公开分享直接下载

    v2 PR1 升级: 当 share 有密码时, 必须 query ?password=xxx 才返流;
    缺少或错密码返 403. 防枚举: 整体返 403 不区分 404 vs 403.
    """
    svc = DriveService(db)
    f = await svc.verify_share_access(token, password=password)
    if f is None:
        # 一律 403 (不区分"不存在/过期/密码错")
        raise HTTPException(status_code=403, detail="分享链接已过期或密码错误")
    if not f.file_path:
        raise HTTPException(status_code=404, detail="分享文件无 MinIO 对象")

    # 公开访问: 不需要 visibility 校验 (token 本身代表 owner 主动授权)
    # 但仍原子 +1 下载计数
    new_count = await svc.increment_download_count(f.id)

    filename = f.file_name or f.title or f"file_{f.id}"
    content_type = f.file_type or "application/octet-stream"
    file_size = await _get_object_size(f.file_path)

    # Range 支持 (与 /files/{id}/download 一致)
    range_header = request.headers.get("Range")
    if range_header:
        m = re.match(r"bytes=(\d+)-(\d*)", range_header)
        if m and file_size is not None:
            start = int(m.group(1))
            end_str = m.group(2)
            if end_str:
                end = min(int(end_str), file_size - 1)
            else:
                end = file_size - 1
            length = end - start + 1
            async def _range_stream():
                chunk = await _download_range(f.file_path, start, length)
                yield chunk
            return StreamingResponse(
                _range_stream(),
                status_code=206,
                media_type=content_type,
                headers={
                    "Content-Range": f"bytes {start}-{end}/{file_size}",
                    "Accept-Ranges": "bytes",
                    "Content-Length": str(length),
                    "Content-Disposition": build_content_disposition("attachment", filename),
                },
            )

    # 完整下载
    headers = {
        "Content-Disposition": build_content_disposition("attachment", filename),
        "Accept-Ranges": "bytes",
    }
    if file_size is not None:
        headers["Content-Length"] = str(file_size)

    async def _full_stream():
        data = await file_service.download_file(f.file_path)
        yield data

    return StreamingResponse(
        _full_stream(),
        media_type=content_type,
        headers=headers,
    )


# ============================================================
# v2 PR4: 文件秒传 (hash) + 版本历史
# ============================================================


class InstantUploadRequest(BaseModel):
    """秒传查询请求 — 前端算完 hash 后 POST"""
    file_hash: str = Field(..., min_length=32, max_length=64, description="MD5/SHA256 hex (32/64 chars)")
    file_name: str = Field(..., min_length=1, max_length=200)
    file_size: int = Field(..., ge=1, le=MAX_DRIVE_FILE_SIZE_BYTES)
    folder_id: Optional[int] = None
    visibility: str = Field("team", pattern="^(private|team|public)$")
    # v2 PR6-P19: 团队共享盘标识 (前端在 specialView='team' 视图上传 = true)
    is_team_shared: bool = Field(False, description="True=团队共享盘上传, 不在个人网盘显示")


class InstantUploadResponse(BaseModel):
    """秒传响应
    - instant=true  命中: 新 file_id 已创建, 不用再上传字节
    - instant=false 未命中: 前端走老 multipart 上传路径
    """
    instant: bool
    file_id: Optional[int] = None
    file_name: Optional[str] = None
    dedup_saved_bytes: int = 0
    file_size: Optional[int] = None
    file_hash: Optional[str] = None
    # miss 时返前端走老路径的提示
    upload_url: Optional[str] = None


class VersionItem(BaseModel):
    """版本历史明细 (一行 = 一次版本)"""
    id: int
    file_id: int
    version_number: int
    file_hash: str
    file_size: int
    uploaded_by: int
    uploaded_by_name: Optional[str] = None
    change_note: Optional[str] = None
    created_at: str
    is_current: bool = False


class RestoreVersionRequest(BaseModel):
    """恢复版本请求 — 可选 change_note"""
    change_note: Optional[str] = Field(None, max_length=500)


@router.post("/files/instant-upload", response_model=InstantUploadResponse)
async def instant_upload(
    body: InstantUploadRequest,
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """秒传 dedup 查询 + 创建 (命中时)

    命中: hash_lookup 找到同 owner 同 hash 文件 → MinIO copy_object 零带宽 → 新 Knowledge 行
    未命中: 返 instant=false, 前端走老 multipart 上传
    """
    svc = DriveService(db)
    try:
        k, saved = await svc.create_instant_upload(
            file_hash=body.file_hash,
            file_name=body.file_name,
            file_size=body.file_size,
            owner_id=user.id,
            folder_id=body.folder_id,
            visibility=body.visibility,
            created_by=user.id,
            is_team_shared=body.is_team_shared,  # v2 PR6-P19
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=e.message)

    if k:
        return InstantUploadResponse(
            instant=True,
            file_id=k.id,
            file_name=k.file_name,
            dedup_saved_bytes=saved,
            file_size=k.file_size,
            file_hash=k.file_hash,
        )
    return InstantUploadResponse(
        instant=False,
        upload_url="/api/v1/drive/files/upload",
    )


@router.get("/files/{file_id}/versions", response_model=List[VersionItem])
async def list_file_versions(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """列文件版本历史 (按 version_number desc)

    权限: 走 _can_see_file, private 文件仅 owner 可看
    """
    svc = DriveService(db)
    try:
        items = await svc.list_versions(file_id=file_id, current_user_id=user.id)
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=e.message)
    return [VersionItem(**item) for item in items]


@router.post("/files/{file_id}/versions/{version_id}/restore", response_model=DriveFileItem)
async def restore_file_version(
    file_id: int,
    version_id: int,
    body: RestoreVersionRequest = RestoreVersionRequest(),
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """恢复历史版本

    创建新行 v{cur.version+1}, file_hash 与历史版一致 (字节级还原)
    旧版本链保留 (cur.is_latest=False), 新行 is_latest=True
    """
    svc = DriveService(db)
    try:
        new_k = await svc.restore_version(
            file_id=file_id,
            version_id=version_id,
            uploader_id=user.id,
            change_note=body.change_note,
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=e.message)
    await db.refresh(new_k)
    return _to_item(new_k)


# ============================================================
# v2 PR5: 配额 + 分片上传 + 断点续传 + 缩略图 (2026-07-01)
# ============================================================


@router.get("/storage-quota", response_model=StorageQuotaResponse)
async def get_storage_quota(
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """v2 PR5: 获取当前用户的网盘配额详情

    返回:
    - used_bytes / quota_bytes / percent: 用于 UI badge 颜色阈值 (≥80% 黄, ≥95% 红)
    - file_count: 活跃 drive 文件数
    - is_over_quota: 已超额 (≤0)
    """
    svc = DriveService(db)
    return await svc.get_storage_quota(user.id)


@router.post("/files/upload/init", response_model=ChunkedUploadInitResponse)
async def init_chunked_upload(
    body: ChunkedUploadInitRequest,
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """v2 PR5: 初始化分片上传 session

    配额检查: 文件大小超配额返 413
    24h TTL: session.expires_at = now + 24h

    返回 upload_id, 前端 chunk 0 写到 PUT /files/upload/{id}/chunk/0
    """
    svc = DriveService(db)
    try:
        session = await svc.init_chunked_upload(
            user_id=user.id,
            file_name=body.file_name,
            file_size=body.file_size,
            total_chunks=body.total_chunks,
            file_hash=body.file_hash,
            folder_id=body.folder_id,
            visibility=body.visibility,
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=e.message)

    return ChunkedUploadInitResponse(
        upload_id=session.id,
        object_name=f"drive-uploads/{session.id}/final",
        total_chunks=session.total_chunks,
        chunk_size_hint=5 * 1024 * 1024,
        uploaded_chunks=list(session.uploaded_chunks or []),
        expires_at=session.expires_at.isoformat(),
    )


@router.put("/files/upload/{upload_id}/chunk/{chunk_index}")
async def upload_chunk(
    upload_id: str,
    chunk_index: int,
    request: Request,
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """v2 PR5: 上传单个 chunk (raw bytes body, NOT multipart)

    - session 不存在 / 已完成 / 已过期 → 404
    - chunk_index 越界 → 400
    - 成功返回 {uploaded_chunks: [0, 1, ...], total_chunks}

    注: 接收 raw bytes 而非 multipart File, 避免 5MB+ chunk 走 multipart 编码膨胀 33%
    """
    chunk_data = await request.body()
    if not chunk_data:
        raise HTTPException(status_code=400, detail="chunk body 为空")

    svc = DriveService(db)
    try:
        session = await svc.upload_chunk(
            session_id=upload_id,
            user_id=user.id,
            chunk_index=chunk_index,
            chunk_data=chunk_data,
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=e.message)

    return {
        "upload_id": upload_id,
        "uploaded_chunks": list(session.uploaded_chunks or []),
        "total_chunks": session.total_chunks,
    }


@router.get("/files/upload/{upload_id}", response_model=ChunkedUploadStatusResponse)
async def get_chunked_upload_status(
    upload_id: str,
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """v2 PR5: 断点续传 - 查 session 已传 chunks 列表

    前端 reload 后调此端点 → 拿到 uploaded_chunks → 跳过这些索引
    """
    svc = DriveService(db)
    session = await svc.get_chunked_session(upload_id, user.id)
    if not session:
        raise HTTPException(status_code=404, detail="Session 不存在或无权访问")

    return ChunkedUploadStatusResponse(
        upload_id=session.id,
        file_name=session.file_name,
        file_size=session.file_size,
        total_chunks=session.total_chunks,
        uploaded_chunks=list(session.uploaded_chunks or []),
        status=session.status,
        expires_at=session.expires_at.isoformat(),
    )


@router.post("/files/upload/{upload_id}/complete", response_model=DriveFileItem)
async def complete_chunked_upload(
    upload_id: str,
    body: ChunkedUploadCompleteRequest = ChunkedUploadCompleteRequest(),
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """v2 PR5: 完成分片上传 → 拼接 → 创建 Knowledge 行

    前置条件: 全部 chunks 已传 (uploaded_chunks == total_chunks)
    副作用:
    - session.status='completed'
    - Fire-and-forget: 重算配额 + 生成缩略图
    - 清 MinIO staging objects
    """
    svc = DriveService(db)
    try:
        new_file = await svc.complete_chunked_upload(
            session_id=upload_id,
            user_id=user.id,
            change_note=body.change_note,
            is_team_shared=body.is_team_shared,  # v2 PR6-P19
        )
    except DriveServiceError as e:
        raise HTTPException(status_code=e.status_code, detail=e.message)

    return _to_item(new_file)


@router.post("/files/upload/{upload_id}/abort", status_code=204)
async def abort_chunked_upload(
    upload_id: str,
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """v2 PR5: 中止分片上传 + 清 MinIO staging"""
    svc = DriveService(db)
    success = await svc.abort_chunked_upload(upload_id, user.id)
    if not success:
        raise HTTPException(status_code=404, detail="Session 不存在或已完成")


@router.get("/files/{file_id}/thumbnail", response_model=ThumbnailResponse)
async def get_thumbnail(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """v2 PR5: 获取文件缩略图信息

    - status=ready → 返 thumbnail_url (MinIO 公开读 URL, 前端 <img src>)
    - status=pending → 返 null URL, 前端 fallback 到 type icon
    - status=failed → 同 pending, UI 可显示 retry 按钮

    越权: 必须能 _can_see_file 才返 (复用 drive_service.get_file 路径)
    """
    svc = DriveService(db)
    k = await svc.get_file(file_id, user.id)
    if not k:
        raise HTTPException(status_code=404, detail="文件不存在或无权访问")

    thumb_url = None
    if k.thumbnail_status == "ready" and k.thumbnail_path:
        # 用 file_service.get_url 返 MinIO 公开读 URL
        thumb_url = file_service.get_url(k.thumbnail_path, expires=3600)

    return ThumbnailResponse(
        file_id=file_id,
        thumbnail_path=k.thumbnail_path,
        thumbnail_status=k.thumbnail_status or "pending",
        thumbnail_url=thumb_url,
    )


# ============================================================
# v2 PR8e: 预览元信息端点 (60s Redis 缓存) — 2026-07-24
# ============================================================
#
# 设计: 与 thumbnail 端点互补
# - thumbnail: 返 200x200 缩略图 (适合 list 卡片, PR5 已有)
# - preview: 返 轻量元信息 (适合 detail 页"封面"块 + 文件选择器)
#   - image: width/height + thumbnail_url
#   - pdf: page_count + first_page_url
#   - text: 前 1KB 文本
#   - other: 仅 file_type/file_size (前端 fallback)
#
# 60s 缓存理由: 预览元信息稳定 (文件不变元信息不变), 但又不能 cache 太久
# (用户上传新版本要能秒看到新页数), 60s 是 mobile feed 类短轮询的甜点
#
# 复用 thumbnail (200x200) 而非新生成 1024px 大图: PR8e 范围严格 1-2 文件,
# 大图渲染留给未来 PR. thumbnail_status=failed 时仍能根据 file_type 给个
# type-icon URL (thumbnail_url=None, 前端用 type icon fallback)


PREVIEW_CACHE_TTL = 60  # Redis TTL (秒)
PREVIEW_TEXT_MAX_BYTES = 1024  # text preview 截断阈值
PREVIEW_IMAGE_EXT = {".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".svg"}
PREVIEW_PDF_EXT = {".pdf"}
PREVIEW_TEXT_EXT = {
    ".txt", ".md", ".json", ".csv", ".tsv", ".xml", ".log",
    ".yaml", ".yml", ".ini", ".conf", ".sh", ".bat",
    ".sql", ".env", ".properties", ".html", ".htm",
}


def _classify_preview_type(file_name: str, file_type: str) -> str:
    """返回 preview_type: image | pdf | text | other

    优先按扩展名判断 (更准), fallback 用 file_type (MIME 字符串).
    """
    import os as _os
    ext = ""
    if file_name:
        ext = _os.path.splitext(file_name)[1].lower()
    if not ext and file_type:
        # file_type 可能是 ".pdf" 或 "application/pdf", 兼容两种
        ft = file_type.lower().strip()
        if ft.startswith("."):
            ext = ft
        elif "/" in ft:
            # 简化映射
            if ft.startswith("image/"):
                return "image"
            if ft == "application/pdf":
                return "pdf"
            if ft.startswith("text/"):
                return "text"
    if ext in PREVIEW_IMAGE_EXT:
        return "image"
    if ext in PREVIEW_PDF_EXT:
        return "pdf"
    if ext in PREVIEW_TEXT_EXT:
        return "text"
    return "other"


async def _read_minio_bytes(object_name: str, max_bytes: int = 5 * 1024 * 1024) -> Optional[bytes]:
    """从 MinIO 读文件字节 (封顶 max_bytes, 5MB 默认, 防止大文件 OOM)

    失败 (文件不存在 / 超限) 返回 None.
    """
    import asyncio as _asyncio
    try:
        def _sync_read() -> bytes:
            response = file_service.client.get_object(
                file_service.bucket, object_name,
            )
            try:
                data = response.read(max_bytes + 1)  # 多读 1 字节判断是否超限
                return data
            finally:
                response.close()
                response.release_conn()
        return await _asyncio.to_thread(_sync_read)
    except Exception as e:
        logger.warning(f"[preview] 读 MinIO 失败 {object_name}: {e}")
        return None


def _probe_image(data: bytes) -> tuple:
    """Pillow 取 image 宽高, 失败返 (None, None).

    性能: 200x200 jpg ~10ms, 5MB png ~50ms (I/O 主导).
    """
    try:
        from PIL import Image
        import io as _io
        img = Image.open(_io.BytesIO(data))
        w, h = img.size
        img.close()
        return w, h
    except Exception:
        return None, None


def _probe_pdf(data: bytes) -> Optional[int]:
    """PyMuPDF 取 pdf 页数, 失败返 None.

    注意: doc.close() 必须 finally, 避免文件句柄泄漏.
    """
    try:
        import fitz  # PyMuPDF
        import io as _io
        doc = fitz.open(stream=data, filetype="pdf")
        try:
            return len(doc)
        finally:
            doc.close()
    except Exception:
        return None


def _read_text_preview(data: bytes) -> str:
    """读前 1KB 文本 (binary 时返空字符串)

    binary 检测: 前 512 字节含 NUL byte 即判 binary.
    UTF-8 容错解码 (fatal=False), 截断时不抛.
    """
    if not data:
        return ""
    # binary 检测
    sample = data[:512]
    for byte in sample:
        if byte == 0:
            return ""
    # UTF-8 解码 + 截断
    preview = data[:PREVIEW_TEXT_MAX_BYTES]
    text = preview.decode("utf-8", errors="replace")
    return text


@router.get("/files/{file_id}/preview", response_model=PreviewResponse)
async def get_file_preview(
    file_id: int,
    db: AsyncSession = Depends(get_db),
    user: Member = Depends(get_current_user),
):
    """v2 PR8e: 获取文件预览元信息 (60s Redis 缓存)

    - image → width/height + thumbnail_url
    - pdf → page_count + first_page_url
    - text → text_preview (前 1KB)
    - other → 仅 file_type/file_size

    越权: 必须能 get_file 才返 (复用 drive_service.get_file)
    失败: MinIO 读失败/Pillow/PyMuPDF 失败 → 字段 None, 不抛 5xx
    """
    import json as _json
    import asyncio as _asyncio
    from app.core.redis import get_redis

    # 1) 越权 + 元信息
    svc = DriveService(db)
    k = await svc.get_file(file_id, user.id)
    if not k:
        raise HTTPException(status_code=404, detail="文件不存在或无权访问")

    # 2) Redis 缓存查询 (per file_id, 含 user_id 防越权穿透)
    cache_key = f"drive:preview:v1:{user.id}:{file_id}"
    redis_client = await get_redis()
    try:
        cached_str = await redis_client.get(cache_key)
        if cached_str:
            cached = _json.loads(cached_str)
            cached["cached"] = True
            return PreviewResponse(**cached)
    except Exception as e:
        # Redis 不可用不阻塞主流程 (降级实时解析)
        logger.warning(f"[preview] Redis get 失败, 降级实时解析: {e}")

    # 3) 实时解析
    file_name = k.file_name or k.title or f"file_{k.id}"
    file_type = k.file_type or ""
    preview_type = _classify_preview_type(file_name, file_type)

    # thumbnail_url 复用 thumbnail_path (MinIO 公开读 URL)
    thumbnail_url = None
    if k.thumbnail_status == "ready" and k.thumbnail_path:
        thumbnail_url = file_service.get_url(k.thumbnail_path, expires=3600)

    width = height = None
    page_count = None
    text_preview = None
    first_page_url = None

    # 仅在有 MinIO 对象时尝试解析元信息
    if k.file_path and preview_type in ("image", "pdf", "text"):
        # 异步读 (限 5MB, 避免大文件 OOM)
        data = await _read_minio_bytes(k.file_path, max_bytes=5 * 1024 * 1024)
        if data is not None:
            if preview_type == "image":
                width, height = await _asyncio.to_thread(_probe_image, data)
            elif preview_type == "pdf":
                page_count = await _asyncio.to_thread(_probe_pdf, data)
                # first page URL 复用 thumbnail_url 模式 (MinIO 公开读)
                # 注意: 这是整个 PDF 的公开读 URL, 浏览器 PDF viewer 默认显示第一页
                if thumbnail_url:
                    first_page_url = thumbnail_url  # 走 thumbnail URL 模式
            elif preview_type == "text":
                text_preview = await _asyncio.to_thread(_read_text_preview, data)

    # 4) 构造响应
    response = PreviewResponse(
        file_id=k.id,
        file_name=file_name,
        file_type=file_type,
        file_size=k.file_size or 0,
        preview_type=preview_type,
        cached=False,
        width=width,
        height=height,
        page_count=page_count,
        thumbnail_url=thumbnail_url,
        text_preview=text_preview,
        first_page_url=first_page_url,
    )

    # 5) 写 Redis 缓存 (best-effort, 失败不阻塞)
    try:
        await redis_client.setex(
            cache_key,
            PREVIEW_CACHE_TTL,
            response.model_dump_json(),
        )
    except Exception as e:
        logger.warning(f"[preview] Redis setex 失败 (非阻塞): {e}")

    return response


# ============================================================
# v2 PR8.5: 移动端聚合 feed 端点 (减少移动端 N 次请求)
# ============================================================

class MobileFeedResponse(BaseModel):
    """PR8.5 mobile feed 响应 (一次返回驱动网盘首页所有数据)
    设计要点:
      - 一次 HTTP 请求 = 减少移动端 N 次往返 (网络延迟对移动端敏感)
      - 各 section 独立 try/except, 失败不阻塞其他 section
      - limit 参数控制每个 section 大小, 默认 10
    """
    recent: List[dict] = Field(default_factory=list)
    starred: List[dict] = Field(default_factory=list)
    team: List[dict] = Field(default_factory=list)
    trash_count: int = 0
    unread_notifications: int = 0
    storage_used_bytes: int = 0
    storage_quota_bytes: int = 0
    generated_at: str = ""


@router.get("/mobile-feed", response_model=MobileFeedResponse)
async def get_mobile_feed(
    limit: int = Query(10, ge=1, le=50, description="每个 section 返回条数"),
    user: Member = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """v2 PR8.5: 移动端首页聚合 (4 sections + 2 stats)

    Sections:
      - recent: 最近修改 drive 文件 (visibility 含自己可见)
      - starred: 用户收藏
      - team: 团队空间最新
      - (server-side) trash_count + unread_notifications

    失败隔离: 任一 section 失败时返空列表, 不抛 5xx 让整个 feed 失败
    """
    from datetime import datetime, timezone
    svc = DriveService(db)
    feed = MobileFeedResponse(
        generated_at=datetime.now(timezone.utc).isoformat(),
    )

    # recent
    try:
        items, _ = await svc.list_files(
            current_user_id=user.id, sort_by="updated_at", sort_order="desc",
            page=1, page_size=limit,
        )
        sids = await svc.get_starred_ids([x.id for x in items], user.id)
        feed.recent = [_drive_file_to_dict(f, user.id, sids) for f in items]
    except Exception as e:
        logger.warning(f"[MobileFeed] recent failed: {e}")

    # starred
    try:
        items, _ = await svc.list_files(
            current_user_id=user.id, sort_by="starred_at", sort_order="desc",
            page=1, page_size=limit, starred_only=True,
        )
        sids = {x.id for x in items}  # starred_only=True 已按本人过滤
        feed.starred = [_drive_file_to_dict(f, user.id, sids) for f in items]
    except Exception as e:
        logger.warning(f"[MobileFeed] starred failed: {e}")

    # team (visibility=team)
    try:
        items, _ = await svc.list_files(
            current_user_id=user.id, sort_by="updated_at", sort_order="desc",
            page=1, page_size=limit, visibility_filter="team",
        )
        sids = await svc.get_starred_ids([x.id for x in items], user.id)
        feed.team = [_drive_file_to_dict(f, user.id, sids) for f in items]
    except Exception as e:
        logger.warning(f"[MobileFeed] team failed: {e}")

    # trash count (PR2 已实现 list_trash)
    try:
        _, total = await svc.list_trash(current_user_id=user.id, page=1, page_size=1)
        feed.trash_count = total
    except Exception as e:
        logger.warning(f"[MobileFeed] trash_count failed: {e}")

    # unread notifications (PR6)
    try:
        from app.services.notification_service import notification_service
        unread = await notification_service.count_unread(db, user_id=user.id)
        feed.unread_notifications = unread
    except Exception as e:
        logger.warning(f"[MobileFeed] unread_notifications failed: {e}")

    # storage stats (PR5 引入 file_size 列后启用 used_bytes/quota_bytes)
    try:
        stats = await svc.storage_stats(user.id)
        # 当前 PR2 阶段 storage_stats 仅返 file_count + by_visibility,
        # used_bytes/quota_bytes 留 0 等 PR5 引入 size 列后填充
        feed.storage_used_bytes = 0
        feed.storage_quota_bytes = 0
    except Exception as e:
        logger.warning(f"[MobileFeed] storage stats failed: {e}")

    return feed


def _drive_file_to_dict(file: Knowledge, user_id: int, starred_ids=None) -> dict:
    """PR8.5 helper: Knowledge → mobile feed dict

    复用 drive_service 已有 _to_dict 模式 (如果有), 这里独立实现避免循环 import

    批次① 收藏个人化: starred_ids (本人 star 集合) 给定时 per-user 覆盖
    is_starred 字段; 未给定时退回 legacy 列 (仅供无 db 上下文的兜底路径)。
    """
    if starred_ids is not None:
        starred_flag = file.id in starred_ids
    else:
        starred_flag = bool(getattr(file, "is_starred", False))
    return {
        "id": file.id,
        "title": file.title,
        "file_name": file.file_name,
        "file_type": file.file_type,
        "file_size": file.file_size,
        "visibility": file.visibility,
        "is_starred": starred_flag,
        "updated_at": file.updated_at.isoformat() if file.updated_at else None,
        "folder_id": getattr(file, "folder_id", None),
    }
