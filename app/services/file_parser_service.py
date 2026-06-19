"""文件内容提取服务 — 支持 PDF/Word/Excel/PPT/TXT/Markdown"""

import asyncio
import io
import logging
import re

logger = logging.getLogger("microbubble.file_parser")


def _find_fig_anchor(text: str, fig_idx: int) -> int:
    """在 text 中找 fig_idx 对应的 "Fig. N" 第一次出现位置

    返回 char offset，失败返回 None

    搜索模式（按优先级）：
    1. "Fig. N" / "Figure N"（精确编号）
    2. "Fig N"（无句点）
    3. 同一 fig_idx 第一次出现即可

    注：fig_idx 1-based（Fig. 1, Fig. 2...）
    """
    if fig_idx < 1:
        return None
    # 多种 Fig 写法
    patterns = [
        rf"\bFig\.?\s*{fig_idx}\b",           # Fig. 1, Fig 1
        rf"\bFigure\s+{fig_idx}\b",            # Figure 1
        rf"\bFig\.\s*{fig_idx}[a-z]\b",        # Fig. 1a
        rf"\bFig\.?\s*{fig_idx}\s*[:.,)]",    # Fig. 1: or Fig. 1.
    ]
    earliest = None
    for pat in patterns:
        m = re.search(pat, text)
        if m:
            if earliest is None or m.start() < earliest:
                earliest = m.end()  # 插在 "Fig. N" 之后
    return earliest


class FileParserService:
    """从各类文件中提取文本和图片"""

    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.md'}

    async def extract_content(self, file_data: bytes, filename: str, content_type: str) -> dict:
        """提取文件文本+图片，返回 {text, images: {placeholder: bytes}}"""
        ext = ''
        if '.' in filename:
            ext = '.' + filename.rsplit('.', 1)[-1].lower()

        if ext == '.pdf' or 'pdf' in content_type:
            return await self._parse_pdf(file_data)
        elif ext in ('.docx',) or 'wordprocessingml' in content_type:
            return {"text": await self._parse_docx(file_data), "images": {}}
        elif ext in ('.xlsx',) or 'spreadsheetml' in content_type:
            return {"text": await self._parse_xlsx(file_data), "images": {}}
        elif ext in ('.pptx',) or 'presentationml' in content_type:
            return {"text": await self._parse_pptx(file_data), "images": {}}
        elif ext in ('.txt', '.md') or content_type.startswith('text/'):
            return {"text": file_data.decode('utf-8', errors='replace'), "images": {}}
        else:
            raise ValueError(f"不支持的文件类型: {ext}")

    async def _parse_pdf(self, data: bytes) -> dict:
        """解析 PDF — 用 PyMuPDF 提取文本和嵌入图片

        2026-06-19 Phase 7 v3: 让 [FIGURE:N] 占位符插入到正文**第一次提到 Fig. N** 的位置
        之前是堆在页尾（PyMuPDF get_text 后插入），导致 inline 后图都堆在页底
        用户读正文时看到 "Fig. 1" 但图片在另一段才出现，体验差

        算法：
        1. 每页提 text + 图片
        2. 解析每页 text 中的 "Fig. N" 引用位置
        3. 对每张图（按 PDF 提取顺序编号 fig_idx）：
           - 在 text 中第一次出现 "Fig. fig_idx" 的位置插入 [FIGURE:fig_idx]
           - 找不到就 fallback 到页尾
        4. 这样 inline 时图就出现在正文中提到它的位置
        """
        def _extract():
            import fitz  # PyMuPDF
            import re

            doc = fitz.open(stream=data, filetype="pdf")
            if doc.needs_pass:
                doc.close()
                raise ValueError("PDF 已加密，无法解析")

            pages_text = []
            images = {}  # placeholder -> image_bytes

            # 跨页跟踪"已用过的 Fig. N 引用位置"，避免重复使用
            # 每张图找到的 anchor 位置（offset in full text），用于后续拼接
            fig_anchor_positions = []  # [(fig_idx, char_position_in_full_text)]

            full_text_parts = []  # 累积 text 用于 anchor 搜索
            current_offset = 0

            for page_num in range(len(doc)):
                page = doc[page_num]

                # [PAGE:N] 标记（用于 inline 定位）
                page_marker = f"[PAGE:{page_num + 1}]"
                pages_text.append(page_marker)
                full_text_parts.append(page_marker)
                current_offset += len(page_marker) + 1  # +1 for \n

                # 提取文本
                text = page.get_text()
                page_text_start_offset = current_offset
                if text:
                    pages_text.append(text)
                    full_text_parts.append(text)
                    current_offset += len(text) + 1

                # 提取嵌入图片
                img_list = page.get_images(full=True)
                for img_idx, img_info in enumerate(img_list):
                    xref = img_info[0]
                    try:
                        base_image = doc.extract_image(xref)
                        img_bytes = base_image["image"]
                        ext_lower = base_image["ext"].lower()
                        if len(img_bytes) < 2048:
                            continue
                    except Exception as e:
                        logger.warning(f"PDF 图片提取失败(page={page_num}, xref={xref}): {e}")
                        continue

                    # fig_idx 编号：基于该页之前已处理的所有图
                    fig_idx = len(fig_anchor_positions) + 1
                    placeholder = f"\n[FIGURE:{fig_idx}]\n"
                    images[placeholder] = {
                        "bytes": img_bytes,
                        "ext": ext_lower if ext_lower in ("png", "jpeg", "jpg") else "png",
                        "page": page_num + 1,
                    }

                    # 关键修复：在**图所在页 + 下一页**找 Fig. fig_idx 引用
                    # 学术 PDF 通常图在上、描述在下（Fig 1 实际渲染在 P3，描述在 P4）
                    # 我们的策略：图 inline 到 Fig. N 引用后，让用户从引用跳到图
                    # 但要把图放对位置（不堆页尾，不放 P1 TOC）
                    local_anchor = None
                    if text:
                        # 先在当前页找
                        local_anchor = _find_fig_anchor(text, fig_idx)

                    if local_anchor is not None:
                        # 当前页找到引用 → 图跟在引用后面
                        global_anchor = page_text_start_offset + local_anchor
                        fig_anchor_positions.append((page_num, global_anchor, placeholder))
                    else:
                        # 当前页没找到 → 放当前页末尾（而不是下页）
                        # 下页引用用户能看到（"Fig. 5"），再点页面里的图查看
                        fig_anchor_positions.append((page_num, current_offset, placeholder))

            doc.close()

            # 现在拼接：在每个 fig_anchor_position 插入 [FIGURE:N] 占位符
            # 用 \n[FIGURE:N]\n 包裹以确保 markdown 渲染时有视觉分隔
            # 按位置倒序插入，避免偏移
            full_text = "".join(full_text_parts) + "\n"  # 末尾加换行

            # 排序：先处理 position 大的（后插入的），再处理小的
            sorted_anchors = sorted(
                [(pos, ph) for _, pos, ph in fig_anchor_positions],
                key=lambda x: x[0],
                reverse=True,
            )
            for pos, placeholder in sorted_anchors:
                # 在 pos 位置插入 placeholder
                # placeholder 格式: \n[FIGURE:N]\n
                full_text = full_text[:pos] + placeholder + full_text[pos:]

            # 用 full_text 替换原来的 pages_text 拼接
            # 注意：pages_text 包含 [PAGE:N] 标记和原 text
            # 我们的 anchor 搜索基于 full_text_parts（包含 [PAGE:N]）
            # 但最终 inline 需要在 text 中能找到 [FIGURE:N]
            # 简化：直接用 full_text 作为最终 text
            text = full_text.strip()
            if not text:
                raise ValueError("PDF 为扫描件或图片型文档，无文字层")
            return {"text": text, "images": images}
        try:
            return await asyncio.to_thread(_extract)
        except ValueError:
            raise
        except Exception as e:
            logger.exception("PDF 解析失败")
            raise ValueError(f"PDF 解析失败: {str(e)}")

    async def _parse_docx(self, data: bytes) -> str:
        """解析 Word 文档"""
        def _extract():
            from docx import Document
            doc = Document(io.BytesIO(data))
            return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        return await asyncio.to_thread(_extract)

    async def _parse_xlsx(self, data: bytes) -> str:
        """解析 Excel 文件"""
        def _extract():
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), read_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    row_text = ' '.join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        texts.append(row_text)
            return '\n'.join(texts)
        return await asyncio.to_thread(_extract)

    async def _parse_pptx(self, data: bytes) -> str:
        """解析 PPT 文件

        2026-06-19 Phase 7 v2: 每页前插入 [PAGE:N] 标记，便于多模态 inline 按页定位
        """
        def _extract():
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            texts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                # 2026-06-19 Phase 7: 每页前加 [PAGE:N] 标记
                texts.append(f"[PAGE:{slide_num}]")
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = ' '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                slide_texts.append(row_text)
                if slide_texts:
                    texts.append(f"--- 第{slide_num}页 ---")
                    texts.extend(slide_texts)
            return '\n'.join(texts)
        return await asyncio.to_thread(_extract)


file_parser_service = FileParserService()
