"""多模态解析服务 — Phase 7

核心职责：
1. 解析 PDF/PPTX 提取嵌入图片（复用 file_parser_service）+ 可选栅格化 PPT 整页
2. 图片上传 MinIO + 写 KnowledgeImage 表
3. 并发 OCR 每张图（OCR 后端由 settings.MULTIMODAL_OCR_BACKEND 选）
4. 写 KnowledgeExtraction 表（公式/表格/图表/OCR 块）
5. 把公式/表格插入到 Knowledge.formatted_content 的合适位置

关键设计：
- 独立容错：每张图 OCR 失败不影响其他图，DB 单图失败 = 1 行 ocr_status=failed
- 并发控制：asyncio.Semaphore(settings.MULTIMODAL_OCR_CONCURRENCY)
- 单文档限流：最多 settings.MULTIMODAL_MAX_IMAGES_PER_DOC 张
- 跳过太小的图（装饰/图标）：< settings.MULTIMODAL_MIN_IMAGE_PIXELS
- 跳过太大：> settings.MULTIMODAL_MAX_IMAGE_PIXELS 等比缩小
- PPT 整页栅格化：可选（settings.MULTIMODAL_RASTERIZE_PPTX，默认 False）
- API 友好：单入口 extract_for_knowledge(knowledge_id) 异步后台执行

调用链：
    KnowledgeService._analyze_and_embed
        → Step 7: multimodal_extraction_service.extract_for_knowledge(knowledge_id)
"""

import asyncio
import io
import logging
import re
from datetime import datetime
from typing import List, Optional, Tuple

from PIL import Image
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.config import settings
from app.core.database import async_session
from app.models.knowledge import Knowledge
from app.models.knowledge_multimodal import (
    KnowledgeImage,
    KnowledgeExtraction,
    EXTRACTION_KIND_FORMULA,
    EXTRACTION_KIND_TABLE,
    EXTRACTION_KIND_CHART,
    EXTRACTION_KIND_IMAGE_BLOCK,
)
from app.services.file_service import file_service
from app.services.ocr_service import ocr_service, OCRBackendError, OCRUnsupportedError, _clean_ocr_text

logger = logging.getLogger("microbubble.multimodal")


# ============================================================================
# 工具函数
# ============================================================================


def _resize_image_if_needed(image_bytes: bytes, max_pixels: int) -> Tuple[bytes, Optional[Tuple[int, int]]]:
    """如果图片像素超过 max_pixels，等比缩小。返回 (bytes, (w, h))"""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        w, h = img.size
        if w * h <= max_pixels:
            return image_bytes, (w, h)
        # 计算缩放比例
        ratio = (max_pixels / (w * h)) ** 0.5
        new_w = int(w * ratio)
        new_h = int(h * ratio)
        # 转换模式（RGBA→RGB for JPEG）
        if img.mode in ("RGBA", "LA", "P") and image_bytes[:2] == b"\xff\xd8":
            img = img.convert("RGB")
        resized = img.resize((new_w, new_h), Image.LANCZOS)
        buf = io.BytesIO()
        # 保持原始格式
        fmt = (img.format or "PNG").upper()
        if fmt == "JPEG" or image_bytes[:2] == b"\xff\xd8":
            resized.save(buf, format="JPEG", quality=85)
            mime = "image/jpeg"
        else:
            resized.save(buf, format="PNG", optimize=True)
            mime = "image/png"
        return buf.getvalue(), (new_w, new_h)
    except Exception as e:
        logger.warning(f"图片缩放失败（用原图）: {e}")
        return image_bytes, None


def _detect_mime(image_bytes: bytes) -> str:
    """从魔数检测 MIME（与 file_service 保持一致的简单版）"""
    if image_bytes[:8] == b"\x89PNG\r\n\x1a\n":
        return "image/png"
    elif image_bytes[:2] == b"\xff\xd8":
        return "image/jpeg"
    elif image_bytes[:4] == b"GIF8":
        return "image/gif"
    elif image_bytes[:4] == b"RIFF" and image_bytes[8:12] == b"WEBP":
        return "image/webp"
    return "image/png"


def _should_skip_image(image_bytes: bytes) -> bool:
    """过滤太小（装饰/图标）的图片"""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        w, h = img.size
        return w * h < settings.MULTIMODAL_MIN_IMAGE_PIXELS
    except Exception:
        return True


# ============================================================================
# PDF / PPTX 解析（提取嵌入图片）
# ============================================================================


def _extract_pdf_images(file_bytes: bytes) -> List[dict]:
    """从 PDF 提取嵌入图片

    Returns: [{bytes, page, ext, width, height}, ...]
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF 未安装，无法提取 PDF 图片")
        return []

    results = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page_num in range(len(doc)):
            page = doc[page_num]
            img_list = page.get_images(full=True)
            for img_info in img_list:
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image["image"]
                    ext = base_image.get("ext", "png").lower()
                    if ext == "jpg":
                        ext = "jpeg"
                    # 过滤小图（< 2KB）
                    if len(img_bytes) < 2048:
                        continue
                    # 获取尺寸
                    try:
                        img = Image.open(io.BytesIO(img_bytes))
                        w, h = img.size
                    except Exception:
                        w = h = None
                    results.append({
                        "bytes": img_bytes,
                        "page": page_num + 1,
                        "ext": ext,
                        "width": w,
                        "height": h,
                    })
                    if len(results) >= settings.MULTIMODAL_MAX_IMAGES_PER_DOC * 2:
                        # 提取时多取一些，OCR 阶段再过滤（< MIN_IMAGE_PIXELS 跳过）
                        break
                except Exception as e:
                    logger.warning(f"PDF 图片提取失败 (page={page_num}, xref={xref}): {e}")
            if len(results) >= settings.MULTIMODAL_MAX_IMAGES_PER_DOC * 2:
                break
        doc.close()
    except Exception as e:
        logger.error(f"PDF 解析失败: {e}")
    return results


def _extract_pptx_images(file_bytes: bytes) -> List[dict]:
    """从 PPTX 提取嵌入图片 + 可选栅格化整页"""
    try:
        from pptx import Presentation
    except ImportError:
        return []

    results = []
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                # 提取嵌入图片
                if shape.shape_type == 13:  # PICTURE
                    try:
                        img_bytes = shape.image.blob
                        ext = shape.image.ext.lower()
                        if ext == "jpg":
                            ext = "jpeg"
                        try:
                            img = Image.open(io.BytesIO(img_bytes))
                            w, h = img.size
                        except Exception:
                            w = h = None
                        results.append({
                            "bytes": img_bytes,
                            "page": slide_num,
                            "ext": ext,
                            "width": w,
                            "height": h,
                        })
                    except Exception as e:
                        logger.warning(f"PPTX 图片提取失败 (slide={slide_num}): {e}")
            if len(results) >= settings.MULTIMODAL_MAX_IMAGES_PER_DOC * 2:
                break
    except Exception as e:
        logger.error(f"PPTX 解析失败: {e}")
    return results


def _extract_file_images(file_bytes: bytes, filename: str, content_type: str) -> List[dict]:
    """根据文件类型分发到对应提取器"""
    ext = ""
    if "." in filename:
        ext = "." + filename.rsplit(".", 1)[-1].lower()
    if ext == ".pdf" or "pdf" in content_type:
        return _extract_pdf_images(file_bytes)
    elif ext == ".pptx" or "presentationml" in content_type:
        return _extract_pptx_images(file_bytes)
    return []


# ============================================================================
# 主入口
# ============================================================================


class MultimodalExtractionService:
    """多模态提取服务

    用法：
        await multimodal_extraction_service.extract_for_knowledge(knowledge_id)
    """

    async def _reset_multimodal_data(self, knowledge_id: int):
        """重跑前清空旧的 image + extraction 记录（防止重复 OCR 数据堆积）

        同时清除 formatted_content / content 中的 MULTIMODAL_INLINED 标记，
        让 inline 重跑会重新生成排版。
        """
        async with async_session() as db:
            # Delete old images (cascade deletes extractions via source_image_id SET NULL)
            await db.execute(
                KnowledgeImage.__table__.delete().where(KnowledgeImage.knowledge_id == knowledge_id)
            )
            # Delete old extractions
            from app.models.knowledge_multimodal import KnowledgeExtraction
            await db.execute(
                KnowledgeExtraction.__table__.delete().where(KnowledgeExtraction.knowledge_id == knowledge_id)
            )
            # Reset analysis_status so UI shows "analyzing" again
            result = await db.execute(
                select(Knowledge).where(Knowledge.id == knowledge_id)
            )
            knowledge = result.scalar_one_or_none()
            if knowledge:
                knowledge.analysis_status = "analyzing"
                # Remove inline marker from content/formatted_content
                if knowledge.content and self.INLINE_MARKER in knowledge.content:
                    knowledge.content = knowledge.content.split(self.INLINE_MARKER)[0].rstrip()
                if knowledge.formatted_content and self.INLINE_MARKER in knowledge.formatted_content:
                    knowledge.formatted_content = knowledge.formatted_content.split(self.INLINE_MARKER)[0].rstrip()
            await db.commit()

    async def extract_for_knowledge(self, knowledge_id: int):
        """异步提取指定知识条目的图片+OCR+公式/表格

        流程：
        0. 清空旧 image + extraction（防止重复）
        1. 取 Knowledge 行（必须有 file_path/file_type）
        2. 从 MinIO 下载原文件
        3. file_parser_service 提取嵌入图片
        4. 过滤/缩放后上传 MinIO + 写 KnowledgeImage
        5. 并发 OCR（classify_and_extract 一次拿 5 字段）
        6. 写 KnowledgeExtraction（去重：同 source_image_id 相似内容跳过）
        7. 把 LaTeX/Markdown table 拼到 formatted_content（可选）
        """
        from app.services.file_service import file_service as fs

        # 0. 清空旧数据（避免重复 OCR 累积）
        await self._reset_multimodal_data(knowledge_id)

        # 1. 取 Knowledge
        async with async_session() as db:
            result = await db.execute(
                select(Knowledge).where(Knowledge.id == knowledge_id)
            )
            knowledge = result.scalar_one_or_none()
            if not knowledge:
                logger.warning(f"knowledge_id={knowledge_id} 不存在")
                return {"ok": False, "reason": "knowledge_not_found"}
            if not knowledge.file_path:
                logger.debug(f"knowledge_id={knowledge_id} 无 file_path 跳过")
                return {"ok": False, "reason": "no_file_path"}
            file_path = knowledge.file_path
            file_name = knowledge.file_name or "unknown"
            file_type = knowledge.file_type or "application/octet-stream"
            title = knowledge.title

        # 2. 下载原文件
        try:
            file_bytes = await fs.download_file(file_path)
        except Exception as e:
            logger.error(f"下载文件失败(knowledge_id={knowledge_id}): {e}")
            return {"ok": False, "reason": f"download_failed: {e}"}

        # 3. 提取图片
        raw_images = _extract_file_images(file_bytes, file_name, file_type)
        if not raw_images:
            logger.info(f"knowledge_id={knowledge_id} 无嵌入图片，跳过多模态")
            return {"ok": True, "skipped": True, "reason": "no_images"}

        # 4. 过滤小图 + 缩放 + 准备上传
        candidates = []
        for img in raw_images:
            if _should_skip_image(img["bytes"]):
                continue
            # 缩放（如果太大）
            resized_bytes, dimensions = _resize_image_if_needed(
                img["bytes"], settings.MULTIMODAL_MAX_IMAGE_PIXELS
            )
            if dimensions and (dimensions[0] and dimensions[1]):
                w, h = dimensions
            else:
                w = img.get("width")
                h = img.get("height")
            candidates.append({
                "bytes": resized_bytes,
                "page": img["page"],
                "ext": img.get("ext", "png"),
                "width": w,
                "height": h,
                "mime": _detect_mime(resized_bytes),
                "size": len(resized_bytes),
            })
            if len(candidates) >= settings.MULTIMODAL_MAX_IMAGES_PER_DOC:
                break

        if not candidates:
            return {"ok": True, "skipped": True, "reason": "all_images_too_small"}

        logger.info(
            f"knowledge_id={knowledge_id}: 准备 OCR {len(candidates)} 张图片"
        )

        # 5. 上传 MinIO + 写 KnowledgeImage（pending 状态）
        image_records = await self._upload_images(knowledge_id, candidates)

        # 6. 并发 OCR
        ocr_results = await self._ocr_images_concurrent(image_records)

        # 7. 写 KnowledgeExtraction
        extraction_counts = await self._save_extractions(knowledge_id, image_records, ocr_results)

        # 8. 把 LaTeX/table 注入 formatted_content
        await self._inject_into_formatted_content(knowledge_id)

        logger.info(
            f"knowledge_id={knowledge_id} 多模态提取完成: "
            f"images={len(image_records)}, "
            f"formulas={extraction_counts['formula']}, "
            f"tables={extraction_counts['table']}, "
            f"charts={extraction_counts['chart']}, "
            f"ocr_blocks={extraction_counts['image_block']}"
        )
        return {
            "ok": True,
            "knowledge_id": knowledge_id,
            "images_total": len(image_records),
            "images_ocr_ok": sum(1 for r in ocr_results if r.get("ok")),
            "extractions": extraction_counts,
        }

    # ── 内部步骤 ───────────────────────────────────────────────────────

    async def _upload_images(
        self, knowledge_id: int, candidates: List[dict]
    ) -> List[KnowledgeImage]:
        """上传图片到 MinIO + 写 KnowledgeImage 行（pending 状态）

        不在这里 commit，分批 commit 减少长事务
        """
        records = []
        async with async_session() as db:
            for idx, cand in enumerate(candidates, 1):
                try:
                    upload = await file_service.upload_file(
                        file_data=cand["bytes"],
                        filename=f"fig_{knowledge_id}_{idx}.{cand['ext']}",
                        content_type=cand["mime"],
                        prefix=f"knowledge/{knowledge_id}/images",
                    )
                    img = KnowledgeImage(
                        knowledge_id=knowledge_id,
                        image_url=upload["url"],
                        image_object_name=upload["object_name"],
                        mime_type=cand["mime"],
                        file_size=cand["size"],
                        page_number=cand["page"],
                        width=cand["width"],
                        height=cand["height"],
                        ocr_status="pending",
                    )
                    db.add(img)
                    await db.flush()  # 拿 id
                    records.append(img)
                except Exception as e:
                    logger.error(f"图片上传失败(knowledge_id={knowledge_id}, idx={idx}): {e}")
            await db.commit()
        return records

    async def _ocr_images_concurrent(
        self, image_records: List[KnowledgeImage]
    ) -> List[dict]:
        """并发 OCR 每张图片（一次 classify_and_extract 拿 5 字段）"""
        sem = ocr_service.semaphore

        async def _process_one(img: KnowledgeImage) -> dict:
            async with sem:
                # 从 MinIO 重新下载（图片对象已在 MinIO）
                try:
                    img_bytes = await file_service.download_file(img.image_object_name)
                except Exception as e:
                    return {"image_id": img.id, "ok": False, "error": f"download: {e}"}
                try:
                    parsed = await ocr_service.classify_and_extract(
                        img_bytes, img.mime_type or "image/png"
                    )
                    return {
                        "image_id": img.id,
                        "ok": True,
                        "parsed": parsed,
                    }
                except (OCRBackendError, OCRUnsupportedError) as e:
                    return {"image_id": img.id, "ok": False, "error": str(e)}
                except Exception as e:
                    return {"image_id": img.id, "ok": False, "error": f"unexpected: {e}"}

        tasks = [_process_one(img) for img in image_records]
        results = await asyncio.gather(*tasks, return_exceptions=False)
        return results

    async def _save_extractions(
        self,
        knowledge_id: int,
        image_records: List[KnowledgeImage],
        ocr_results: List[dict],
    ) -> dict:
        """把 OCR 结果写入 KnowledgeExtraction + 更新 KnowledgeImage.ocr_status

        重要：image_records 来自 _upload_images 的另一个 session，
        不能直接 mutate（SQLAlchemy session 隔离，mutate 不会持久化）。
        必须在新 session 里 re-fetch 后再更新。

        去重策略（防 OCR 重复入库）：
        - 同一 OCR 调用只算一个 result，不会内部重复
        - 不同图之间如果 chart_description 高度相似（前 200 字相同）跳过
        - image_block 完全相同的 text 跳过
        """
        from sqlalchemy.orm.attributes import flag_modified
        counts = {"formula": 0, "table": 0, "chart": 0, "image_block": 0, "dedup_skipped": 0}
        # 索引 image_id → page_number/position_data（从前一个 session 的快照取）
        meta_by_id = {img.id: {
            "page_number": img.page_number,
            "position_data": img.position_data,
        } for img in image_records}

        # 去重跟踪：跨图去重（同一文档内）
        seen_chart_descs: set = set()  # 200 字前缀
        seen_image_blocks: set = set()  # 完整内容

        async with async_session() as db:
            # 一次性 re-fetch 所有 image（用新 session 绑定）
            img_ids = [r["image_id"] for r in ocr_results]
            if img_ids:
                fetched = await db.execute(
                    select(KnowledgeImage).where(KnowledgeImage.id.in_(img_ids))
                )
                img_by_id = {img.id: img for img in fetched.scalars().all()}
            else:
                img_by_id = {}

            for result in ocr_results:
                img_id = result["image_id"]
                img = img_by_id.get(img_id)
                if not img:
                    # image 在新 session 找不到（极少见）
                    continue

                if not result.get("ok"):
                    img.ocr_status = "failed"
                    img.ocr_error = (result.get("error") or "")[:1000]
                    continue

                parsed = result.get("parsed", {})
                category = parsed.get("category", "figure")
                ocr_text = parsed.get("text", "") or ""
                latex = parsed.get("latex")
                table_md = parsed.get("table_md")
                chart_desc = parsed.get("chart_description")
                caption = parsed.get("caption")
                model_used = settings.VISION_MODEL

                # 主 OCR 文本（清洗 thinking 残留）
                if ocr_text:
                    img.ocr_text = _clean_ocr_text(ocr_text)[:10000]
                if caption and not ocr_text:
                    img.ocr_text = f"[caption] {_clean_ocr_text(caption)[:500]}"
                img.ocr_status = "done"
                img.ocr_model = model_used
                img.ocr_at = datetime.utcnow()

                # 公式提取
                if latex and latex.strip() and latex.strip() != "NO_FORMULA":
                    ext = KnowledgeExtraction(
                        knowledge_id=knowledge_id,
                        source_image_id=img_id,
                        kind=EXTRACTION_KIND_FORMULA,
                        page_number=meta_by_id.get(img_id, {}).get("page_number"),
                        position_data=meta_by_id.get(img_id, {}).get("position_data"),
                        data={"latex": latex.strip(), "caption": caption},
                        content_text=latex.strip(),
                        confidence=0.8,
                        model_used=model_used,
                        source="auto",
                    )
                    db.add(ext)
                    counts["formula"] += 1

                # 表格提取
                if table_md and table_md.strip() and table_md.strip() != "NO_TABLE":
                    parsed_table = _parse_markdown_table(table_md)
                    ext = KnowledgeExtraction(
                        knowledge_id=knowledge_id,
                        source_image_id=img_id,
                        kind=EXTRACTION_KIND_TABLE,
                        page_number=meta_by_id.get(img_id, {}).get("page_number"),
                        position_data=meta_by_id.get(img_id, {}).get("position_data"),
                        data=parsed_table,
                        content_text=table_md.strip(),
                        confidence=0.8,
                        model_used=model_used,
                        source="auto",
                    )
                    db.add(ext)
                    counts["table"] += 1

                # 图表描述（去重：跨图相似内容跳过）
                if chart_desc and chart_desc.strip():
                    desc = chart_desc.strip()
                    # 用前 200 字做去重指纹（避免 LLM 输出末尾差异）
                    fingerprint = desc[:200]
                    if fingerprint in seen_chart_descs:
                        counts["dedup_skipped"] += 1
                    else:
                        seen_chart_descs.add(fingerprint)
                        ext = KnowledgeExtraction(
                            knowledge_id=knowledge_id,
                            source_image_id=img_id,
                            kind=EXTRACTION_KIND_CHART,
                            page_number=meta_by_id.get(img_id, {}).get("page_number"),
                            position_data=meta_by_id.get(img_id, {}).get("position_data"),
                            data={"description": desc},
                            content_text=desc,
                            confidence=0.7,
                            model_used=model_used,
                            source="auto",
                        )
                        db.add(ext)
                        counts["chart"] += 1

                # image_block 总是记（即使没特殊提取，也有 OCR 文字）
                # 去重：完全相同的 text 跳过（防同图重复入库）
                if ocr_text:
                    cleaned_text = _clean_ocr_text(ocr_text)
                    if cleaned_text in seen_image_blocks:
                        counts["dedup_skipped"] += 1
                    else:
                        seen_image_blocks.add(cleaned_text)
                        ext = KnowledgeExtraction(
                            knowledge_id=knowledge_id,
                            source_image_id=img_id,
                            kind=EXTRACTION_KIND_IMAGE_BLOCK,
                            page_number=meta_by_id.get(img_id, {}).get("page_number"),
                            position_data=meta_by_id.get(img_id, {}).get("position_data"),
                            data={"text": cleaned_text, "caption": caption},
                            content_text=cleaned_text,
                            confidence=0.6,
                            model_used=model_used,
                            source="auto",
                        )
                        db.add(ext)
                        counts["image_block"] += 1

            await db.commit()
        return counts

    async def _inject_into_formatted_content(self, knowledge_id: int):
        """把多模态提取的公式/表格追加到 formatted_content 末尾

        不破坏原有内容，只在末尾追加"## 多模态提取"section
        """
        async with async_session() as db:
            result = await db.execute(
                select(Knowledge).where(Knowledge.id == knowledge_id)
            )
            knowledge = result.scalar_one_or_none()
            if not knowledge:
                return

            # 取所有 active 提取物
            ext_result = await db.execute(
                select(KnowledgeExtraction)
                .where(
                    KnowledgeExtraction.knowledge_id == knowledge_id,
                    KnowledgeExtraction.is_active == True,
                )
                .order_by(KnowledgeExtraction.kind, KnowledgeExtraction.id)
            )
            extractions = ext_result.scalars().all()

            formulas = [e for e in extractions if e.kind == EXTRACTION_KIND_FORMULA]
            tables = [e for e in extractions if e.kind == EXTRACTION_KIND_TABLE]
            charts = [e for e in extractions if e.kind == EXTRACTION_KIND_CHART]

            if not (formulas or tables or charts):
                return

            sections = []
            if formulas:
                sections.append("### 提取的公式\n")
                for f in formulas:
                    latex = (f.data or {}).get("latex", "")
                    if latex:
                        sections.append(f"$$\n{latex}\n$$")
                        if f.content_text and f.content_text != latex:
                            sections.append(f"<!-- {f.content_text} -->")
            if tables:
                sections.append("### 提取的表格\n")
                for t in tables:
                    md = t.content_text or ""
                    if md:
                        sections.append(md)
                        sections.append("")
            if charts:
                sections.append("### 图表说明\n")
                for c in charts:
                    desc = (c.data or {}).get("description", "")
                    if desc:
                        sections.append(f"- {desc}")

            injection = "\n\n## 多模态提取\n\n" + "\n\n".join(sections) + "\n"

            base = knowledge.formatted_content or knowledge.content or ""
            # 避免重复注入
            if "## 多模态提取" in base:
                return
            knowledge.formatted_content = base + injection
            await db.commit()

    # ── Inline placement（v2: 把图/表/公式 嵌回原文位置） ─────────

    INLINE_MARKER = "<!-- MULTIMODAL_INLINED v2 -->"
    PAGE_MARKER_RE = re.compile(r"\[PAGE:(\d+)\]")
    FIGURE_MARKER_RE = re.compile(r"\[FIGURE:(\d+)\]")
    # 装饰图过滤：标题/版权/封面/目录/logo 等首页元素
    DECORATIVE_KEYWORDS = (
        "elsevier", "journal of", "journal homepage", "www.elsevier",
        "highlights", "highlights", "article info", "abstract",
        "graphical abstract", "keywords:",
    )

    @classmethod
    def _is_decorative_image(cls, img) -> bool:
        """判断一张图是否为装饰图（journal logo / 封面 / 目录等）

        装饰图特征：
        - OCR 文字里有 ELSEVIER / Journal name / Highlights / Article info 等关键词
        - OCR 文字太短（< 50 字符，无实质内容）
        - 标题类文本
        """
        ocr = (img.ocr_text or "").lower()
        if not ocr:
            return True
        if len(ocr) < 50:
            return True
        for kw in cls.DECORATIVE_KEYWORDS:
            if kw in ocr:
                return True
        return False

    async def reparse_pdf_with_page_markers(self, knowledge_id: int) -> dict:
        """重新解析 PDF 加上 [PAGE:N] 标记，更新 knowledge.content

        适用场景：知识入库时是 Phase 7 v1 之前，没有 page marker → inline 无法定位
        """
        from app.services.file_service import file_service as fs
        from app.services.file_parser_service import file_parser_service

        async with async_session() as db:
            result = await db.execute(
                select(Knowledge).where(Knowledge.id == knowledge_id)
            )
            knowledge = result.scalar_one_or_none()
            if not knowledge:
                return {"ok": False, "reason": "knowledge_not_found"}
            if not knowledge.file_path:
                return {"ok": False, "reason": "no_file_path"}
            file_path = knowledge.file_path
            file_name = knowledge.file_name or "unknown"
            file_type = knowledge.file_type or "application/octet-stream"

        try:
            file_data = await fs.download_file(file_path)
        except Exception as e:
            return {"ok": False, "reason": f"download_failed: {e}"}

        try:
            extracted = await file_parser_service.extract_content(
                file_data, file_name, file_type
            )
        except Exception as e:
            return {"ok": False, "reason": f"parse_failed: {e}"}

        new_content = extracted["text"]
        async with async_session() as db:
            r = await db.execute(select(Knowledge).where(Knowledge.id == knowledge_id))
            k = r.scalar_one_or_none()
            if k:
                k.content = new_content
                # 保留现有 formatted_content 不变（AI 排版有章节结构）
                await db.commit()
        return {"ok": True, "new_content_length": len(new_content), "page_markers": len(self.PAGE_MARKER_RE.findall(new_content))}

    async def inline_extractions_to_content(self, knowledge_id: int) -> dict:
        """用 [FIGURE:N] 占位符位置把多模态提取物嵌回原文对应位置

        设计目标：
        - 解决用户痛点「所有图都在底下，正文只看到 Fig.5 文字」
        - 不用 page_number（首页 8 张装饰图会堆一起），用 [FIGURE:N] 占位符
        - file_parser_service._parse_pdf 在 PDF 提图时，会在 text 里插 [FIGURE:N] 标记
          这些标记就在 caption 旁边 = 原排版位置
        - 流程：
          1. 找 content 里所有 [FIGURE:N] 标记
          2. 对每张图，按 file_parser 提取时的 N 找到对应占位符
          3. 用 markdown inline 替换占位符
        - 装饰图（journal logo / 封面 / 目录）过滤掉
        - 找不到占位符 / 是装饰图 → 退回末尾
        - 幂等：检查 INLINE_MARKER 避免重复插入
        """
        # 1. 拿 knowledge
        async with async_session() as db:
            result = await db.execute(
                select(Knowledge).where(Knowledge.id == knowledge_id)
            )
            knowledge = result.scalar_one_or_none()
            if not knowledge:
                return {"ok": False, "reason": "knowledge_not_found"}
            base_content = knowledge.content or ""

        # 2. 拿所有 images + extractions
        async with async_session() as db:
            img_result = await db.execute(
                select(KnowledgeImage)
                .where(KnowledgeImage.knowledge_id == knowledge_id)
                .order_by(KnowledgeImage.id)
            )
            images = img_result.scalars().all()
            ext_result = await db.execute(
                select(KnowledgeExtraction)
                .where(
                    KnowledgeExtraction.knowledge_id == knowledge_id,
                    KnowledgeExtraction.is_active == True,
                )
                .order_by(KnowledgeExtraction.id)
            )
            extractions = ext_result.scalars().all()

        # 过滤装饰图：ELSEVIER / journal / highlights / 太短 OCR 等
        # 同时排除 OCR 文字少于 30 字符的图（icon / 装饰）
        figure_images = [
            img for img in images
            if (img.ocr_text and len(img.ocr_text) > 30) and not self._is_decorative_image(img)
        ]
        inline_extractions = [
            e for e in extractions
            if e.kind in (EXTRACTION_KIND_FORMULA, EXTRACTION_KIND_TABLE, EXTRACTION_KIND_CHART)
        ]

        if not (figure_images or inline_extractions):
            return {"ok": True, "skipped": True, "reason": "no_inline_candidates"}

        # 3. 找所有 [FIGURE:N] 占位符位置
        figure_positions: dict = {}  # fig_idx (int) -> position (start of [FIGURE:N])
        for m in self.FIGURE_MARKER_RE.finditer(base_content):
            fig_idx = int(m.group(1))
            if fig_idx not in figure_positions:
                figure_positions[fig_idx] = m.start()

        if not figure_positions:
            logger.warning(
                f"knowledge_id={knowledge_id} content 无 [FIGURE:N] 占位符，"
                f"请先调 reparse_pdf_with_page_markers 重解析"
            )
            return {
                "ok": False,
                "reason": "no_figure_placeholders",
                "hint": "POST /knowledge/{id}/reparse-pdf first",
            }

        # 4. 对每张图，找它对应的 [FIGURE:N] 占位符
        #    图片按 id 顺序对应 [FIGURE:1], [FIGURE:2], ...
        #    跳过已用作装饰图的（用 id 的"figure_index"或 fallback 顺序）
        match_positions: list = []  # [(pos, mid, kind, item)]
        skipped_decorative = 0

        # 顺序匹配：knowledge_images 按 id 升序 → 对应 figure_positions[1], [2], ...
        # 但因为装饰图被过滤掉，需要紧凑匹配
        # 策略：图按 id 升序，依次取下一个可用的 [FIGURE:N] 位置
        fig_idx_iter = iter(sorted(figure_positions.keys()))
        available_fig_indices = list(fig_idx_iter)
        img_idx = 0

        for img in images:  # 用全部 images（包括装饰）按 id 顺序对应
            # 跳过装饰图
            if not (img.ocr_text and len(img.ocr_text) > 30) or self._is_decorative_image(img):
                # 占位符也要"消耗"一个（保持位置对齐）
                if available_fig_indices:
                    available_fig_indices.pop(0)
                skipped_decorative += 1
                continue
            if not available_fig_indices:
                # 占位符用完了
                break
            fig_idx = available_fig_indices.pop(0)
            pos = figure_positions[fig_idx]
            match_positions.append((pos, img.id, "image", img))
            img_idx += 1

        # 5. 公式 / 表格 / 图表：page_marker inline（这些在 PDF 里没有占位符）
        #    退一步用 page-based 定位
        page_marker_positions: dict = {}
        for m in self.PAGE_MARKER_RE.finditer(base_content):
            page_num = int(m.group(1))
            if page_num not in page_marker_positions:
                page_marker_positions[page_num] = m.end()

        for e in inline_extractions:
            page = e.page_number or 1
            if page_marker_positions:
                max_page = max(page_marker_positions.keys())
                if page > max_page:
                    page = max_page
                pos = page_marker_positions.get(page)
                if pos is not None:
                    match_positions.append((pos, e.id, "extraction", e))

        if not match_positions:
            return {"ok": True, "skipped": True, "reason": "no_matches"}

        # 6. 按位置从后往前排序（避免 offset 漂移）
        match_positions.sort(key=lambda x: (x[0], -x[1]), reverse=True)

        new_content = base_content
        inserted_mids: set = set()

        # 7. 对每个 [FIGURE:N] 位置，inline 替换
        #    关键：把 [FIGURE:N] 占位符本身替换为 markdown
        #    （不是插入在后面，是替换——这样图就出现在原始 PDF 排版位置）
        for pos, mid, kind, item in match_positions:
            if mid in inserted_mids:
                continue
            if kind != "image":
                # extractions 用 page-based，继续用"插入"模式
                inline_md = self._format_inline_markdown(kind, item)
                if inline_md:
                    new_content = new_content[:pos] + "\n\n" + inline_md + "\n\n" + new_content[pos:]
                    inserted_mids.add(mid)
                    # 偏移调整
                    inserted_len = len(inline_md) + 6
                    for i, m in enumerate(match_positions):
                        if m[1] in inserted_mids:
                            continue
                        if m[0] > pos:
                            match_positions[i] = (m[0] + inserted_len, m[1], m[2], m[3])
                continue

            # image: 替换 [FIGURE:N] 占位符为 inline markdown
            inline_md = self._format_inline_markdown(kind, item)
            if not inline_md:
                continue
            # 占位符长度：[FIGURE:N] 加可能的 \n
            placeholder_end = base_content.find("\n", pos) if pos >= 0 else pos
            if placeholder_end == -1 or placeholder_end > pos + 20:
                placeholder_end = pos + base_content[pos:].find("]") + 1
            else:
                placeholder_end += 1  # 包含 \n
            # 实际占位符 [FIGURE:N]\n 大约 12 字符
            placeholder_text = base_content[pos:pos + 12]  # 截取占位符
            # 安全：找到 [FIGURE:N] 结束位置（] 后到下一个 \n 或 12 字符内）
            end_match = self.FIGURE_MARKER_RE.match(base_content[pos:pos + 15])
            if end_match:
                placeholder_end = pos + end_match.end()
                # 跳过 \n
                if base_content[placeholder_end:placeholder_end + 1] == "\n":
                    placeholder_end += 1
            new_content = new_content[:pos] + inline_md + "\n" + new_content[placeholder_end:]
            inserted_mids.add(mid)
            # 后续 item 的 pos 调整（因为 markdown 比占位符长）
            inserted_len = len(inline_md) + 1 - (placeholder_end - pos)
            if inserted_len != 0:
                for i, m in enumerate(match_positions):
                    if m[1] in inserted_mids:
                        continue
                    if m[0] > pos:
                        match_positions[i] = (m[0] + inserted_len, m[1], m[2], m[3])

        # 8. unmatched items（没占位符的）→ 末尾 append
        item_lookup = {img.id: ("image", img) for img in figure_images}
        for e in inline_extractions:
            item_lookup[e.id] = ("extraction", e)
        unmatched_to_append = []
        for pos, mid, kind, item in match_positions:
            if mid in inserted_mids or mid not in item_lookup:
                continue
            inline_md = self._format_inline_markdown(kind, item)
            if inline_md:
                unmatched_to_append.append(inline_md)

        # 9. 幂等检查
        if self.INLINE_MARKER in new_content:
            new_content = new_content.split(self.INLINE_MARKER)[0].rstrip()

        if unmatched_to_append:
            new_content += "\n\n---\n\n## 未在正文匹配的多模态提取\n\n" + "\n\n".join(unmatched_to_append) + "\n"

        new_content += f"\n\n{self.INLINE_MARKER}\n"

        # 10. 写回
        async with async_session() as db:
            result = await db.execute(
                select(Knowledge).where(Knowledge.id == knowledge_id)
            )
            knowledge = result.scalar_one_or_none()
            if knowledge:
                knowledge.content = new_content
                knowledge.formatted_content = new_content
                await db.commit()

        logger.info(
            f"inline_extractions(knowledge_id={knowledge_id}): "
            f"matched={len(inserted_mids)} skipped_decorative={skipped_decorative} "
            f"unmatched={len(unmatched_to_append)}"
        )
        return {
            "ok": True,
            "knowledge_id": knowledge_id,
            "matches_total": len(inserted_mids),
            "unmatched_total": len(unmatched_to_append),
            "skipped_decorative": skipped_decorative,
            "new_content_length": len(new_content),
        }

    def _format_inline_markdown(self, kind: str, item) -> Optional[str]:
        """把图片 / 公式 / 表格 / 图表 格式化成 inline markdown 片段

        Returns: markdown 字符串（不含前后空行）或 None（无法格式化）
        """
        if kind == "image":
            # 图片：![caption](url)
            img = item
            page = img.page_number
            ocr = (img.ocr_text or "").strip().replace("\n", " ")[:100]
            cap = f"图（P{page}，{ocr}...）" if page else f"图（{ocr}...）"
            return f"![{cap}]({img.image_url})"
        elif kind == "extraction":
            e = item
            if e.kind == EXTRACTION_KIND_FORMULA:
                latex = (e.data or {}).get("latex", "").strip() or e.content_text or ""
                if not latex:
                    return None
                page = e.page_number
                return f"$$\n{latex}\n$$\n<!-- 公式（P{page}） -->" if page else f"$$\n{latex}\n$$"
            elif e.kind == EXTRACTION_KIND_TABLE:
                md = e.content_text or ""
                if not md:
                    return None
                page = e.page_number
                cap = f"**表格（P{page}）**\n\n" if page else "**表格**\n\n"
                return cap + md
            elif e.kind == EXTRACTION_KIND_CHART:
                desc = (e.data or {}).get("description", "").strip() or e.content_text or ""
                if not desc:
                    return None
                page = e.page_number
                header = f"> 📊 **图表说明（P{page}）**" if page else "> 📊 **图表说明**"
                # 多行描述加引用块
                return f"{header}\n> {desc}"
        return None


def _parse_markdown_table(md: str) -> dict:
    """简单解析 Markdown 表格为 {headers, rows, caption}

    容错：只解析标准的 | col1 | col2 | 格式
    """
    if not md:
        return {"headers": [], "rows": [], "caption": None}
    lines = [l.strip() for l in md.split("\n") if l.strip()]
    # 提取 caption（以 ** 开头且不含 |）
    caption = None
    if lines and lines[0].startswith("**") and "|" not in lines[0]:
        caption = lines[0].strip("*").strip()
        lines = lines[1:]
    if not lines:
        return {"headers": [], "rows": [], "caption": caption}

    # 第一行：表头
    def _parse_row(line: str) -> list:
        s = line.strip().strip("|")
        return [c.strip() for c in s.split("|")]

    headers = _parse_row(lines[0])
    rows = []
    for line in lines[1:]:
        if re.match(r"^\|?[\s\-:|]+\|?$", line):
            continue  # 跳过分隔行
        if "|" in line:
            rows.append(_parse_row(line))
    return {"headers": headers, "rows": rows, "caption": caption}


# 全局单例
multimodal_extraction_service = MultimodalExtractionService()
